"""
RFID 설정 시스템 생산 JIG 견적서 PPT 생성 스크립트
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor
from pptx.oxml.xmlchemy import OxmlElement
from pptx.oxml.ns import qn
import os

# 색상 정의
DARK_BG = RGBColor(26, 26, 46)  # #1a1a2e
PURPLE_PRIMARY = RGBColor(102, 126, 234)  # #667eea
PURPLE_SECONDARY = RGBColor(118, 75, 162)  # #764ba2
WHITE = RGBColor(255, 255, 255)
LIGHT_GRAY = RGBColor(170, 170, 170)
GREEN = RGBColor(74, 222, 128)  # #4ade80
DARK_CELL = RGBColor(30, 30, 50)
DARK_CELL2 = RGBColor(40, 40, 70)
DARK_CELL3 = RGBColor(50, 50, 80)
GRAY_TEXT = RGBColor(100, 100, 100)

def set_slide_background(slide, color):
    """슬라이드 배경색 설정"""
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_text_box(slide, left, top, width, height, text, font_size=18, font_color=WHITE, bold=False, align=PP_ALIGN.LEFT):
    """텍스트 박스 추가"""
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = font_color
    p.font.bold = bold
    p.alignment = align
    return txBox

def add_shape_with_text(slide, left, top, width, height, text, shape_type=MSO_SHAPE.ROUNDED_RECTANGLE,
                        fill_color=PURPLE_PRIMARY, text_color=WHITE, font_size=14):
    """도형과 텍스트 추가"""
    shape = slide.shapes.add_shape(shape_type, Inches(left), Inches(top), Inches(width), Inches(height))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()

    tf = shape.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = text_color
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER

    return shape

def create_presentation():
    prs = Presentation()
    prs.slide_width = Inches(13.333)  # 16:9 비율
    prs.slide_height = Inches(7.5)

    # ========== 슬라이드 1: 표지 ==========
    slide1 = prs.slides.add_slide(prs.slide_layouts[6])  # 빈 레이아웃
    set_slide_background(slide1, DARK_BG)

    # UTTEC 로고
    add_text_box(slide1, 0, 1.5, 13.333, 0.8, "UTTEC", font_size=48, font_color=PURPLE_PRIMARY, bold=True, align=PP_ALIGN.CENTER)

    # 제목
    add_text_box(slide1, 0, 2.5, 13.333, 1.2, "RFID 설정 시스템\n생산 JIG 개발", font_size=44, font_color=WHITE, bold=True, align=PP_ALIGN.CENTER)

    # 부제목
    add_text_box(slide1, 0, 4.0, 13.333, 0.6, "견적서", font_size=28, font_color=PURPLE_PRIMARY, bold=False, align=PP_ALIGN.CENTER)

    # 메타 정보
    add_text_box(slide1, 0, 5.5, 13.333, 0.5, "2025년 12월 28일  |  To: RfTech  |  견적 유효기간: 30일",
                 font_size=16, font_color=LIGHT_GRAY, align=PP_ALIGN.CENTER)

    # 슬라이드 번호
    add_text_box(slide1, 12, 7, 1, 0.3, "1 / 6", font_size=12, font_color=LIGHT_GRAY, align=PP_ALIGN.RIGHT)

    # ========== 슬라이드 2: 프로젝트 개요 ==========
    slide2 = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide2, DARK_BG)

    # 제목
    add_text_box(slide2, 0, 0.5, 13.333, 0.8, "프로젝트 개요", font_size=36, font_color=PURPLE_PRIMARY, bold=True, align=PP_ALIGN.CENTER)

    # 3개의 카드
    card_y = 1.8
    card_width = 3.5
    gap = 0.5
    start_x = (13.333 - (card_width * 3 + gap * 2)) / 2

    # Notebook UI
    add_shape_with_text(slide2, start_x, card_y, card_width, 0.8, "Notebook UI", font_size=18)
    add_text_box(slide2, start_x, card_y + 1.0, card_width, 2.5,
                 "생산자용 GUI 프로그램\nWindows 10/11 지원\n직관적인 조작 인터페이스",
                 font_size=14, font_color=LIGHT_GRAY, align=PP_ALIGN.CENTER)

    # ESP32 Dongle
    add_shape_with_text(slide2, start_x + card_width + gap, card_y, card_width, 0.8, "ESP32 Dongle", font_size=18)
    add_text_box(slide2, start_x + card_width + gap, card_y + 1.0, card_width, 2.5,
                 "RFID Reader 내장\nUSB-UART 변환\n프로토콜 중계",
                 font_size=14, font_color=LIGHT_GRAY, align=PP_ALIGN.CENTER)

    # 타겟 시스템
    add_shape_with_text(slide2, start_x + (card_width + gap) * 2, card_y, card_width, 0.8, "타겟 시스템 설정", font_size=18)
    add_text_box(slide2, start_x + (card_width + gap) * 2, card_y + 1.0, card_width, 2.5,
                 "RFID Tag 2개 등록\n평일/휴일 근무시간\nRTC 시간 동기화",
                 font_size=14, font_color=LIGHT_GRAY, align=PP_ALIGN.CENTER)

    # 시스템 흐름도
    flow_y = 5.5
    add_shape_with_text(slide2, 1.5, flow_y, 2.5, 0.7, "Notebook", font_size=14)
    add_text_box(slide2, 4.1, flow_y + 0.15, 0.8, 0.4, "USB", font_size=12, font_color=LIGHT_GRAY, align=PP_ALIGN.CENTER)
    add_shape_with_text(slide2, 5, flow_y, 2.5, 0.7, "ESP32 Dongle", font_size=14)
    add_text_box(slide2, 7.6, flow_y + 0.15, 0.8, 0.4, "UART", font_size=12, font_color=LIGHT_GRAY, align=PP_ALIGN.CENTER)
    add_shape_with_text(slide2, 8.5, flow_y, 2.5, 0.7, "타겟 시스템", font_size=14, fill_color=PURPLE_SECONDARY)

    add_text_box(slide2, 12, 7, 1, 0.3, "2 / 6", font_size=12, font_color=LIGHT_GRAY, align=PP_ALIGN.RIGHT)

    # ========== 슬라이드 3: 견적 내역 ==========
    slide3 = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide3, DARK_BG)

    add_text_box(slide3, 0, 0.3, 13.333, 0.7, "견적 내역", font_size=36, font_color=PURPLE_PRIMARY, bold=True, align=PP_ALIGN.CENTER)

    # 테이블 데이터
    table_data = [
        ["항목", "수량", "단가 (원)", "금액 (원)"],
        ["[하드웨어]", "", "", ""],
        ["ESP32 Dongle (RFID Reader, 케이스 포함)", "2 SET", "150,000", "300,000"],
        ["UART 연결 케이블 (3.3V TTL, 1m)", "2 SET", "10,000", "20,000"],
        ["USB 케이블 (USB-A to Micro-USB)", "2 SET", "5,000", "10,000"],
        ["하드웨어 소계", "", "", "330,000"],
        ["[소프트웨어 개발]", "", "", ""],
        ["ESP32 펌웨어 개발", "1 식", "800,000", "800,000"],
        ["Notebook UI 프로그램 개발", "1 식", "1,200,000", "1,200,000"],
        ["통신 프로토콜 설계 및 구현", "1 식", "300,000", "300,000"],
        ["시스템 통합 테스트", "1 식", "200,000", "200,000"],
        ["소프트웨어 소계", "", "", "2,500,000"],
        ["[문서 및 기타]", "", "", ""],
        ["기술 문서 작성", "1 식", "200,000", "200,000"],
        ["문서/기타 소계", "", "", "200,000"],
        ["합계 (VAT 별도)", "", "", "3,030,000"],
    ]

    # 테이블 생성
    rows = len(table_data)
    cols = 4
    table_left = Inches(1)
    table_top = Inches(1.1)
    table_width = Inches(11.333)
    table_height = Inches(5.8)

    table = slide3.shapes.add_table(rows, cols, table_left, table_top, table_width, table_height).table

    # 열 너비 설정
    table.columns[0].width = Inches(5.5)
    table.columns[1].width = Inches(1.5)
    table.columns[2].width = Inches(2)
    table.columns[3].width = Inches(2.333)

    # 테이블 데이터 채우기
    for i, row_data in enumerate(table_data):
        for j, cell_text in enumerate(row_data):
            cell = table.cell(i, j)
            cell.text = cell_text

            # 셀 텍스트 서식
            paragraph = cell.text_frame.paragraphs[0]
            paragraph.font.size = Pt(11)
            paragraph.font.color.rgb = WHITE

            # 헤더 행
            if i == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = PURPLE_PRIMARY
                paragraph.font.bold = True
                paragraph.alignment = PP_ALIGN.CENTER
            # 카테고리 행
            elif row_data[0].startswith("["):
                cell.fill.solid()
                cell.fill.fore_color.rgb = DARK_CELL2
                paragraph.font.color.rgb = PURPLE_PRIMARY
                paragraph.font.bold = True
            # 소계 행
            elif "소계" in row_data[0]:
                cell.fill.solid()
                cell.fill.fore_color.rgb = DARK_CELL3
                paragraph.font.bold = True
            # 합계 행
            elif "합계" in row_data[0]:
                cell.fill.solid()
                cell.fill.fore_color.rgb = PURPLE_PRIMARY
                paragraph.font.bold = True
                paragraph.font.size = Pt(13)
            else:
                cell.fill.solid()
                cell.fill.fore_color.rgb = DARK_CELL

            # 금액 열 정렬
            if j >= 2 and i > 0:
                paragraph.alignment = PP_ALIGN.RIGHT
                if j == 3 and cell_text:
                    paragraph.font.color.rgb = GREEN

    add_text_box(slide3, 12, 7, 1, 0.3, "3 / 6", font_size=12, font_color=LIGHT_GRAY, align=PP_ALIGN.RIGHT)

    # ========== 슬라이드 4: 납품 상세 ==========
    slide4 = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide4, DARK_BG)

    add_text_box(slide4, 0, 0.3, 13.333, 0.7, "납품 상세", font_size=36, font_color=PURPLE_PRIMARY, bold=True, align=PP_ALIGN.CENTER)

    # 4개의 상자
    box_width = 5.8

    # 하드웨어 납품물
    add_shape_with_text(slide4, 0.8, 1.2, box_width, 0.5, "하드웨어 납품물", font_size=16)
    add_text_box(slide4, 0.8, 1.8, box_width, 2.0,
                 "• ESP32 Dongle 완제품 2세트\n• RFID Reader 모듈 내장 (RC522)\n• 플라스틱 케이스\n• UART/USB 케이블",
                 font_size=12, font_color=LIGHT_GRAY)

    # 소프트웨어 납품물
    add_shape_with_text(slide4, 6.8, 1.2, box_width, 0.5, "소프트웨어 납품물", font_size=16)
    add_text_box(slide4, 6.8, 1.8, box_width, 2.0,
                 "• Notebook UI 프로그램 (설치 패키지)\n• ESP32 펌웨어 (바이너리 + 소스코드)\n• UI 프로그램 소스코드\n• 빌드 환경 구축 가이드",
                 font_size=12, font_color=LIGHT_GRAY)

    # 기술 문서
    add_shape_with_text(slide4, 0.8, 4.0, box_width, 0.5, "기술 문서", font_size=16)
    add_text_box(slide4, 0.8, 4.6, box_width, 2.0,
                 "• 사용자 매뉴얼 (생산자용)\n• 통신 프로토콜 명세서\n• 하드웨어 회로도\n• 트러블슈팅 가이드",
                 font_size=12, font_color=LIGHT_GRAY)

    # 기술 지원
    add_shape_with_text(slide4, 6.8, 4.0, box_width, 0.5, "기술 지원", font_size=16)
    add_text_box(slide4, 6.8, 4.6, box_width, 2.0,
                 "• 무상 기술 지원: 납품 후 3개월\n• 버그 수정 무상 지원\n• 원격 기술 상담\n• 추가 개발: 별도 협의",
                 font_size=12, font_color=LIGHT_GRAY)

    add_text_box(slide4, 12, 7, 1, 0.3, "4 / 6", font_size=12, font_color=LIGHT_GRAY, align=PP_ALIGN.RIGHT)

    # ========== 슬라이드 5: 견적 요약 ==========
    slide5 = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide5, DARK_BG)

    add_text_box(slide5, 0, 0.5, 13.333, 0.7, "견적 요약", font_size=36, font_color=PURPLE_PRIMARY, bold=True, align=PP_ALIGN.CENTER)

    # 중앙 박스
    box = slide5.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(2.5), Inches(1.5), Inches(8.333), Inches(4.5))
    box.fill.solid()
    box.fill.fore_color.rgb = DARK_CELL2
    box.line.color.rgb = PURPLE_PRIMARY
    box.line.width = Pt(2)

    add_text_box(slide5, 2.5, 1.8, 8.333, 0.5, "총 견적 금액", font_size=20, font_color=LIGHT_GRAY, align=PP_ALIGN.CENTER)
    add_text_box(slide5, 2.5, 2.5, 8.333, 1.0, "₩ 3,030,000", font_size=56, font_color=PURPLE_PRIMARY, bold=True, align=PP_ALIGN.CENTER)
    add_text_box(slide5, 2.5, 3.5, 8.333, 0.4, "* 부가가치세(VAT) 별도", font_size=14, font_color=LIGHT_GRAY, align=PP_ALIGN.CENTER)

    # 조건 정보
    conditions_y = 4.3
    add_text_box(slide5, 3, conditions_y, 2, 0.3, "견적 유효기간", font_size=12, font_color=LIGHT_GRAY, align=PP_ALIGN.CENTER)
    add_text_box(slide5, 3, conditions_y + 0.3, 2, 0.4, "30일", font_size=16, font_color=WHITE, bold=True, align=PP_ALIGN.CENTER)

    add_text_box(slide5, 5.2, conditions_y, 2, 0.3, "납품 예정", font_size=12, font_color=LIGHT_GRAY, align=PP_ALIGN.CENTER)
    add_text_box(slide5, 5.2, conditions_y + 0.3, 2, 0.4, "계약 후 4주", font_size=16, font_color=WHITE, bold=True, align=PP_ALIGN.CENTER)

    add_text_box(slide5, 7.4, conditions_y, 2.5, 0.3, "결제 조건", font_size=12, font_color=LIGHT_GRAY, align=PP_ALIGN.CENTER)
    add_text_box(slide5, 7.4, conditions_y + 0.3, 2.5, 0.4, "계약금 50% / 잔금 50%", font_size=14, font_color=WHITE, bold=True, align=PP_ALIGN.CENTER)

    add_text_box(slide5, 12, 7, 1, 0.3, "5 / 6", font_size=12, font_color=LIGHT_GRAY, align=PP_ALIGN.RIGHT)

    # ========== 슬라이드 6: 연락처 ==========
    slide6 = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide6, DARK_BG)

    add_text_box(slide6, 0, 1.5, 13.333, 1.0, "감사합니다", font_size=48, font_color=PURPLE_PRIMARY, bold=True, align=PP_ALIGN.CENTER)

    add_text_box(slide6, 0, 2.8, 13.333, 0.8, "본 견적서에 대한 문의사항이 있으시면\n아래 연락처로 문의해 주시기 바랍니다.",
                 font_size=18, font_color=LIGHT_GRAY, align=PP_ALIGN.CENTER)

    # 연락처 정보
    contact_y = 4.2
    add_text_box(slide6, 3.5, contact_y, 2, 0.3, "회사명", font_size=12, font_color=LIGHT_GRAY, align=PP_ALIGN.CENTER)
    add_text_box(slide6, 3.5, contact_y + 0.3, 2, 0.5, "UTTEC", font_size=20, font_color=WHITE, bold=True, align=PP_ALIGN.CENTER)

    add_text_box(slide6, 5.7, contact_y, 2, 0.3, "전화", font_size=12, font_color=LIGHT_GRAY, align=PP_ALIGN.CENTER)
    add_text_box(slide6, 5.7, contact_y + 0.3, 2, 0.5, "010-3922-1809", font_size=18, font_color=WHITE, bold=True, align=PP_ALIGN.CENTER)

    add_text_box(slide6, 7.9, contact_y, 2.5, 0.3, "이메일", font_size=12, font_color=LIGHT_GRAY, align=PP_ALIGN.CENTER)
    add_text_box(slide6, 7.9, contact_y + 0.3, 2.5, 0.5, "uttec@uttec.co.kr", font_size=18, font_color=WHITE, bold=True, align=PP_ALIGN.CENTER)

    add_text_box(slide6, 0, 5.8, 13.333, 0.8, "본 견적서는 RfTech 귀사에 제출하는 공식 견적서입니다.\n견적 내용은 상호 협의를 통해 조정될 수 있습니다.",
                 font_size=14, font_color=GRAY_TEXT, align=PP_ALIGN.CENTER)

    add_text_box(slide6, 12, 7, 1, 0.3, "6 / 6", font_size=12, font_color=LIGHT_GRAY, align=PP_ALIGN.RIGHT)

    # 저장
    output_path = os.path.join(os.path.dirname(__file__), "견적서_RFID설정시스템_JIG.pptx")
    prs.save(output_path)
    print(f"PPT 파일 생성 완료: {output_path}")
    return output_path

if __name__ == "__main__":
    create_presentation()

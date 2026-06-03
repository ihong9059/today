"""동아정밀공업 미팅용 PPT (12장) 생성
- 1920x1080 (16:9 widescreen)
- 한국어 / 비즈니스 톤
- 각 슬라이드 notes에 narration
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn
from copy import deepcopy

# ============================================================
# Brand palette
# ============================================================
NAVY = RGBColor(0x10, 0x2A, 0x4C)      # primary dark
NAVY_LIGHT = RGBColor(0x1E, 0x3A, 0x5F)
ORANGE = RGBColor(0xF1, 0x7A, 0x21)    # accent (UTTEC orange)
GOLD = RGBColor(0xC9, 0xA2, 0x27)
GREEN_OK = RGBColor(0x2E, 0x8B, 0x57)
RED_FOCUS = RGBColor(0xC0, 0x39, 0x2B)
CREAM = RGBColor(0xFA, 0xF6, 0xEC)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
GRAY_TEXT = RGBColor(0x33, 0x33, 0x33)
GRAY_LIGHT = RGBColor(0xE0, 0xE0, 0xE0)
GRAY_MID = RGBColor(0x88, 0x88, 0x88)

FONT_KR = "맑은 고딕"  # Malgun Gothic — widely available on Windows

# ============================================================
# Setup presentation (16:9, 1920x1080)
# ============================================================
prs = Presentation()
prs.slide_width = Inches(13.333)  # 1920 / 144 dpi
prs.slide_height = Inches(7.5)    # 1080 / 144 dpi

blank_layout = prs.slide_layouts[6]  # blank

# ============================================================
# Helpers
# ============================================================
def add_rect(slide, x, y, w, h, fill=NAVY, line=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    if line is None:
        shape.line.fill.background()
    else:
        shape.line.color.rgb = line
    shape.shadow.inherit = False
    return shape

def add_text(slide, x, y, w, h, text, size=18, bold=False, color=GRAY_TEXT,
             align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, font=FONT_KR):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.05)
    tf.margin_right = Inches(0.05)
    tf.margin_top = Inches(0.02)
    tf.margin_bottom = Inches(0.02)
    tf.vertical_anchor = anchor
    lines = text.split("\n") if isinstance(text, str) else text
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        run = p.add_run()
        run.text = line
        run.font.name = font
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.color.rgb = color
    return tb

def add_notes(slide, narration):
    notes_tf = slide.notes_slide.notes_text_frame
    notes_tf.text = narration

def add_page_header(slide, slide_no, total, section_label):
    """top thin bar + page indicator"""
    # top accent bar
    add_rect(slide, Inches(0), Inches(0), prs.slide_width, Inches(0.12), fill=ORANGE)
    # logo / brand top-left
    add_text(slide, Inches(0.4), Inches(0.22), Inches(3.0), Inches(0.4),
             "UTTEC", size=14, bold=True, color=NAVY)
    add_text(slide, Inches(0.4), Inches(0.5), Inches(5.0), Inches(0.35),
             "동아정밀공업 귀하 — PET 두께 측정기 2호기 제안", size=10, color=GRAY_MID)
    # page no top-right
    add_text(slide, Inches(11.8), Inches(0.28), Inches(1.3), Inches(0.4),
             f"{slide_no:02d} / {total:02d}  ·  {section_label}",
             size=9, color=GRAY_MID, align=PP_ALIGN.RIGHT)

def add_page_footer(slide):
    add_rect(slide, Inches(0), Inches(7.35), prs.slide_width, Inches(0.15),
             fill=NAVY)
    add_text(slide, Inches(0.4), Inches(7.37), Inches(8), Inches(0.13),
             "UTTEC · 홍광선  ihong9059@gmail.com · 010-7186-2452 · 경기 용인시 기흥구 흥덕유타워",
             size=8, color=WHITE)
    add_text(slide, Inches(11.5), Inches(7.37), Inches(1.6), Inches(0.13),
             "2026-06-03",
             size=8, color=WHITE, align=PP_ALIGN.RIGHT)

def add_section_title(slide, title, subtitle=None):
    add_text(slide, Inches(0.6), Inches(0.95), Inches(12), Inches(0.7),
             title, size=32, bold=True, color=NAVY)
    if subtitle:
        add_text(slide, Inches(0.6), Inches(1.62), Inches(12), Inches(0.4),
                 subtitle, size=15, color=GRAY_MID)
    # accent underline
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                  Inches(0.6), Inches(1.55), Inches(0.8), Inches(0.05))
    line.fill.solid()
    line.fill.fore_color.rgb = ORANGE
    line.line.fill.background()

def add_chip(slide, x, y, w, h, text, fill, text_color=WHITE, size=11, bold=True):
    chip = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    chip.fill.solid()
    chip.fill.fore_color.rgb = fill
    chip.line.fill.background()
    tf = chip.text_frame
    tf.margin_left = Inches(0.08)
    tf.margin_right = Inches(0.08)
    tf.margin_top = Inches(0.02)
    tf.margin_bottom = Inches(0.02)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = text
    r.font.name = FONT_KR
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.color.rgb = text_color
    return chip

def add_card(slide, x, y, w, h, fill=WHITE, border=GRAY_LIGHT, shadow_navy=False):
    if shadow_navy:
        # subtle shadow rectangle behind
        add_rect(slide, x + Inches(0.05), y + Inches(0.05), w, h, fill=GRAY_LIGHT)
    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    card.fill.solid()
    card.fill.fore_color.rgb = fill
    card.line.color.rgb = border
    card.line.width = Pt(0.75)
    card.shadow.inherit = False
    return card

# ============================================================
# Slide 1 — 타이틀
# ============================================================
TOTAL = 10
narration_s1 = """안녕하십니까. 동아정밀공업 관계자 여러분. 본 영상은 서보 기반 비접촉식 PET 용기 두께 측정기 2호기 제어 시스템 구축 제안을 위해 유티텍이 준비한 발표 자료입니다. 함께해 주셔서 감사드립니다."""

s = prs.slides.add_slide(blank_layout)
# full navy background
add_rect(s, Inches(0), Inches(0), prs.slide_width, prs.slide_height, fill=NAVY)
# orange accent band
add_rect(s, Inches(0), Inches(4.2), prs.slide_width, Inches(0.08), fill=ORANGE)

add_text(s, Inches(0.6), Inches(0.6), Inches(6), Inches(0.5),
         "UTTEC  ·  유티텍", size=14, bold=True, color=ORANGE)

add_text(s, Inches(0.6), Inches(1.5), Inches(12), Inches(1.1),
         "동아정밀공업 귀하", size=28, bold=False, color=CREAM)
add_text(s, Inches(0.6), Inches(2.05), Inches(12.4), Inches(0.85),
         "서보 기반 비접촉식 PET 용기",
         size=38, bold=True, color=WHITE)
add_text(s, Inches(0.6), Inches(2.75), Inches(12.4), Inches(0.85),
         "두께 측정기 (2호기)",
         size=38, bold=True, color=WHITE)
add_text(s, Inches(0.6), Inches(3.55), Inches(12), Inches(0.55),
         "PLC + SCADA 제어 시스템 구축 제안",
         size=22, bold=False, color=CREAM)

# divider
add_text(s, Inches(0.6), Inches(4.7), Inches(12), Inches(0.5),
         "발표 :  홍광선 (유티텍 대표, 펌웨어 40년)  ·  임호균 (공장자동화 38년)",
         size=16, color=WHITE)
add_text(s, Inches(0.6), Inches(5.2), Inches(12), Inches(0.4),
         "프로젝트 :  위시캣 #155220   ·   기간 90일   ·   인건비 20,000,000원 (자재 별도)",
         size=14, color=GRAY_LIGHT)

# bottom signature
add_text(s, Inches(0.6), Inches(6.6), Inches(12), Inches(0.4),
         "2026-06-03",
         size=12, color=GRAY_LIGHT)
add_text(s, Inches(0.6), Inches(6.95), Inches(12), Inches(0.35),
         "ihong9059@gmail.com   ·   경기 용인시 기흥구 흥덕유타워",
         size=10, color=GRAY_MID)

add_notes(s, narration_s1)

# ============================================================
# Slide 2 — 인사 + 팀 소개
# ============================================================
narration_s2 = """발표는 유티텍 대표 홍광선과 회로·공장자동화 협업자 임호균이 진행합니다. 두 사람의 합산 경력은 78년이며, 펌웨어 40년, 회로 25년, 공장자동화 38년의 깊이를 한 팀으로 묶어, 회로 설계부터 펌웨어, PLC 제어 로직, SCADA 운영, 데이터 로깅까지 단일 책임으로 수행하는 구조입니다. 양산 제품 6종이 현재 판매 중이며 KC, 일본 TELEC, 유럽 CE 3개국 인증을 보유하고 있습니다."""

s = prs.slides.add_slide(blank_layout)
add_page_header(s, 2, TOTAL, "팀 소개")
add_page_footer(s)
add_section_title(s, "한 팀, 단일 책임", "펌웨어 40년 + 회로/공장자동화 38년 = 합산 78년")

# Two cards
card_w = Inches(5.8)
card_h = Inches(3.6)
card_y = Inches(2.3)

# Card 1 — 홍광선
add_card(s, Inches(0.6), card_y, card_w, card_h, fill=WHITE)
add_rect(s, Inches(0.6), card_y, card_w, Inches(0.5), fill=NAVY)
add_text(s, Inches(0.8), card_y + Inches(0.08), card_w, Inches(0.4),
         "홍광선  ·  유티텍 대표",
         size=18, bold=True, color=WHITE)
add_text(s, Inches(0.8), card_y + Inches(0.7), card_w - Inches(0.4), Inches(0.45),
         "펌웨어 · PLC · SCADA · 데이터 인프라",
         size=15, bold=True, color=ORANGE)
bullets_h = """• 40년 펌웨어 경력 (삼성전자 AV·STB, 파나소닉 LCD/Printer)
• STM32 양산 다수 (F756/F407) + KC 인증
• RPi CM4 EtherCAT MFC 산업 양산
• FastAPI + PostgreSQL + TimescaleDB 백엔드
• SCADA 5요소 자체 구현 라이브 운영"""
add_text(s, Inches(0.85), card_y + Inches(1.25), card_w - Inches(0.5), Inches(2.3),
         bullets_h, size=12, color=GRAY_TEXT)

# Card 2 — 임호균
add_card(s, Inches(6.9), card_y, card_w, card_h, fill=WHITE)
add_rect(s, Inches(6.9), card_y, card_w, Inches(0.5), fill=NAVY)
add_text(s, Inches(7.1), card_y + Inches(0.08), card_w, Inches(0.4),
         "임호균  ·  공장자동화 협업자",
         size=18, bold=True, color=WHITE)
add_text(s, Inches(7.1), card_y + Inches(0.7), card_w - Inches(0.4), Inches(0.45),
         "회로 · 전장 배선 · 모션 · 패널 설계",
         size=15, bold=True, color=ORANGE)
bullets_l = """• 38년 공장자동화 + 모터 제어 경력
• 삼성전자 Video / Motor Controller Board 설계
• AuthenTec(美) ARM CPU + ASIC 설계
• QuickLogic(美) FPGA 설계
• 대한전선 WINDER 권취기 (CAN/RS485 + 모션 양산)"""
add_text(s, Inches(7.15), card_y + Inches(1.25), card_w - Inches(0.5), Inches(2.3),
         bullets_l, size=12, color=GRAY_TEXT)

# bottom: 인증·양산 chips
chip_y = Inches(6.3)
add_chip(s, Inches(0.6), chip_y, Inches(2.1), Inches(0.45),
         "양산 제품 6종 판매 중", NAVY, WHITE, size=12)
add_chip(s, Inches(2.85), chip_y, Inches(1.8), Inches(0.45),
         "KC 인증 (한국)", GREEN_OK, WHITE, size=12)
add_chip(s, Inches(4.8), chip_y, Inches(1.9), Inches(0.45),
         "TELEC 인증 (일본)", GREEN_OK, WHITE, size=12)
add_chip(s, Inches(6.85), chip_y, Inches(1.8), Inches(0.45),
         "CE 인증 (유럽)", GREEN_OK, WHITE, size=12)
add_chip(s, Inches(8.8), chip_y, Inches(3.7), Inches(0.45),
         "일본 BLE Mesh 3,800대 수출 양산", ORANGE, WHITE, size=12)

add_notes(s, narration_s2)

# ============================================================
# Slide 3 — 본 프로젝트 이해 1호기 → 2호기
# ============================================================
narration_s3 = """본 프로젝트의 본질은 1호기 PC 기반 제어의 한계를 산업 표준 PLC와 SCADA로 전환하여 2호기를 전면 재개발하는 사업입니다. 1호기는 OS와 드라이버 의존성, 단일 PC 장애, 유지보수 확장 어려움의 구조적 한계를 안고 있습니다. 2호기는 네 가지 가치로 전환합니다. 첫째, PLC의 결정론적 사이클과 전원 견고성으로 24시간 무인 운영 안정성을 확보합니다. 둘째, SCADA 터치 HMI로 비숙련 작업자도 레시피 선택과 알람 인지를 직관적으로 수행할 수 있습니다. 셋째, 제품 종류별 파라미터 셋을 통합 데이터베이스에 저장·호출하여 품종 변경을 1-Click으로 전환합니다. 넷째, SCADA 표준 히스토리언과 시계열 데이터베이스로 측정 데이터의 트레이스빌리티를 확보합니다."""

s = prs.slides.add_slide(blank_layout)
add_page_header(s, 3, TOTAL, "본 프로젝트 이해")
add_page_footer(s)
add_section_title(s, "1호기 PC → 2호기 PLC + SCADA",
                  "산업 표준 전환으로 얻는 4가지 핵심 가치")

# Left card: 1호기 한계
left_x = Inches(0.6)
right_x = Inches(7.1)
col_w = Inches(5.6)
col_y = Inches(2.3)
col_h = Inches(4.6)

add_card(s, left_x, col_y, col_w, col_h, fill=RGBColor(0xFA, 0xEC, 0xEC))
add_rect(s, left_x, col_y, col_w, Inches(0.5), fill=RED_FOCUS)
add_text(s, left_x + Inches(0.2), col_y + Inches(0.08), col_w, Inches(0.4),
         "1호기  ·  PC (C#) 기반 — 구조적 한계",
         size=16, bold=True, color=WHITE)
limits = """✗  OS / 드라이버 의존성 — Windows 업데이트 영향
✗  단일 PC 장애 위험 — 라인 정지
✗  유지보수 / 기능 확장 어려움
✗  산업 EMI 환경 미보장 — 신뢰성 격차
✗  표준 Historian / Alarm / Recipe 부재"""
add_text(s, left_x + Inches(0.3), col_y + Inches(0.85), col_w - Inches(0.4), Inches(3.5),
         limits, size=14, color=GRAY_TEXT)

# Right card: 2호기 4가지 가치
add_card(s, right_x, col_y, col_w, col_h, fill=RGBColor(0xEC, 0xF3, 0xEA))
add_rect(s, right_x, col_y, col_w, Inches(0.5), fill=GREEN_OK)
add_text(s, right_x + Inches(0.2), col_y + Inches(0.08), col_w, Inches(0.4),
         "2호기  ·  PLC + SCADA — 4가지 핵심 가치",
         size=16, bold=True, color=WHITE)
values = """① 산업 안정성  —  결정론적 사이클 / 24/7 무인 운영
② 운영자 편의성  —  터치 HMI / Recipe 1-Click 전환
③ Recipe 관리  —  품종별 파라미터 셋 통합 DB
④ 데이터 로깅 신뢰성  —  Historian + 시계열 DB"""
add_text(s, right_x + Inches(0.3), col_y + Inches(0.85), col_w - Inches(0.4), Inches(3.5),
         values, size=14, color=GRAY_TEXT)

add_notes(s, narration_s3)

# ============================================================
# Slide 4 — 턴키 역량 (매니저 1순위)
# ============================================================
narration_s4 = """본 프로젝트가 분업으로 진행되면 회로사·펌웨어사·PLC사·SCADA사 사이의 인터페이스 정의와 일정 동기, 책임 경계가 새로운 리스크로 등장합니다. 유티텍은 이 다섯 단계를 한 팀 두 사람이 단일 책임으로 수행합니다. 회로 설계와 전장 패널은 회로 25년 경력의 임호균이, 펌웨어와 PLC 로직, SCADA 운영, 데이터 인프라는 펌웨어 40년 경력의 홍광선이 맡습니다. 현장에서 회로 이슈나 모션 이슈가 발생해도 외부 협력사 호출 없이 즉시 자체 해결이 가능합니다. 이것이 분업 리스크 없는 진정한 턴키 구조입니다."""

s = prs.slides.add_slide(blank_layout)
add_page_header(s, 4, TOTAL, "턴키 역량")
add_page_footer(s)
add_section_title(s, "전장 배선 → PLC → SCADA — 한 팀 책임",
                  "분업 리스크 없는 단일 책임 구조 (매니저 요구 1순위)")

# 5-step chain
chain_y = Inches(3.0)
chain_h = Inches(1.2)
step_w = Inches(2.3)
gap = Inches(0.18)
steps = [
    ("①  회로 / 패널", "임호균", NAVY),
    ("②  펌웨어", "홍광선", NAVY_LIGHT),
    ("③  PLC 로직", "홍광선", NAVY_LIGHT),
    ("④  SCADA 운영", "홍광선", NAVY_LIGHT),
    ("⑤  데이터 로깅", "홍광선", NAVY),
]
start_x = Inches(0.45)
for i, (label, who, color) in enumerate(steps):
    x = start_x + i * (step_w + gap)
    add_card(s, x, chain_y, step_w, chain_h, fill=WHITE)
    add_rect(s, x, chain_y, step_w, Inches(0.45), fill=color)
    add_text(s, x + Inches(0.1), chain_y + Inches(0.06), step_w, Inches(0.35),
             label, size=14, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text(s, x, chain_y + Inches(0.55), step_w, Inches(0.4),
             who, size=15, bold=True, color=NAVY, align=PP_ALIGN.CENTER)
    add_text(s, x, chain_y + Inches(0.88), step_w, Inches(0.3),
             "단일 책임", size=10, color=GRAY_MID, align=PP_ALIGN.CENTER)
    if i < len(steps) - 1:
        # arrow
        arrow_x = x + step_w + Inches(0.01)
        arr = s.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW,
                                 arrow_x, chain_y + Inches(0.45),
                                 gap - Inches(0.02), Inches(0.3))
        arr.fill.solid()
        arr.fill.fore_color.rgb = ORANGE
        arr.line.fill.background()

# bottom message
add_text(s, Inches(0.6), Inches(5.0), Inches(12), Inches(0.5),
         "외부 협력사 호출 없음  ·  인터페이스 갈등 없음  ·  일정 동기 갈등 없음",
         size=18, bold=True, color=NAVY, align=PP_ALIGN.CENTER)
add_text(s, Inches(0.6), Inches(5.6), Inches(12), Inches(0.45),
         "현장 회로 / 모션 이슈 발생 시 외주 호출 없이 즉시 자체 해결",
         size=14, color=GRAY_TEXT, align=PP_ALIGN.CENTER)

# Comparison badges
add_chip(s, Inches(2.5), Inches(6.4), Inches(3.8), Inches(0.5),
         "분업 — 인터페이스 리스크", RED_FOCUS, WHITE, size=13)
add_chip(s, Inches(7.0), Inches(6.4), Inches(3.8), Inches(0.5),
         "턴키 — 단일 책임 (UTTEC)", GREEN_OK, WHITE, size=13)

add_notes(s, narration_s4)

# ============================================================
# Slide 5 — 양산 자산 6종 + 인증
# ============================================================
narration_s5 = """저희가 양산 출하 중인 제품을 보여드립니다. STM32F756 기반 RS-485 Modbus RTU 컴프레서 밸브 컨트롤러로 KC 인증을 보유하고 있고, STM32F407 세탁기 컨트롤러, 라즈베리파이 CM4 EtherCAT MFC 컨트롤러, 라즈베리파이 3 V-Cut 컨트롤러로 중국 수출 실적, nRF52832 BLE 온도 컨트롤러, 그리고 일본 AMANO 자전거주차장 LED Dimmer 3,800대 BLE Mesh 양산까지 모두 현재 판매 중입니다. 산업 안정성을 입증하는 KC, 일본 TELEC, 유럽 CE 3개국 인증을 보유하고 있어, 본 프로젝트의 산업 EMI·EMC 환경 대응을 동일한 패턴으로 즉시 적용할 수 있습니다."""

s = prs.slides.add_slide(blank_layout)
add_page_header(s, 5, TOTAL, "양산 자산")
add_page_footer(s)
add_section_title(s, "양산 제품 6종 + 3개국 인증",
                  "현재 판매 중인 산업 양산 자산")

products = [
    ("STM32F756", "컴프레서 밸브 컨트롤러", "RS-485 Modbus RTU · FreeRTOS · PID", "KC 인증", ORANGE),
    ("STM32F407", "세탁기 컨트롤러", "FreeRTOS · 모터 시퀀스 양산", "양산 판매중", NAVY),
    ("RPi CM4", "EtherCAT MFC 컨트롤러", "산업 결정론 사이클 · 정밀 제어", "양산 판매중", NAVY),
    ("RPi 3", "V-Cut 컨트롤러", "정밀 위치 이송", "중국 수출", GOLD),
    ("nRF52832", "BLE 온도 컨트롤러", "저전력 BLE · 양산", "양산 판매중", NAVY),
    ("AMANO BLE Mesh", "자전거주차장 LED Dimmer", "외부안테나 200m+ · BLE Mesh", "일본 3,800대 수출", ORANGE),
]
cols = 3
rows = 2
cell_w = Inches(4.05)
cell_h = Inches(2.0)
grid_x = Inches(0.6)
grid_y = Inches(2.3)
gap_x = Inches(0.22)
gap_y = Inches(0.22)
for i, (model, name, desc, badge, badge_color) in enumerate(products):
    r = i // cols
    c = i % cols
    x = grid_x + c * (cell_w + gap_x)
    y = grid_y + r * (cell_h + gap_y)
    add_card(s, x, y, cell_w, cell_h, fill=WHITE)
    # left accent strip
    add_rect(s, x, y, Inches(0.12), cell_h, fill=NAVY)
    add_text(s, x + Inches(0.25), y + Inches(0.12), cell_w - Inches(0.4), Inches(0.4),
             model, size=14, bold=True, color=ORANGE)
    add_text(s, x + Inches(0.25), y + Inches(0.5), cell_w - Inches(0.4), Inches(0.4),
             name, size=14, bold=True, color=NAVY)
    add_text(s, x + Inches(0.25), y + Inches(0.95), cell_w - Inches(0.4), Inches(0.5),
             desc, size=11, color=GRAY_TEXT)
    add_chip(s, x + Inches(0.25), y + Inches(1.55), Inches(2.5), Inches(0.35),
             badge, badge_color, WHITE, size=10)

# Bottom certifications
add_text(s, Inches(0.6), Inches(6.55), Inches(12), Inches(0.4),
         "3개국 산업 인증 보유  →  본 프로젝트 EMI/EMC 양산 패턴 즉시 적용",
         size=14, bold=True, color=NAVY, align=PP_ALIGN.CENTER)

add_notes(s, narration_s5)

# ============================================================
# Slide 7 — 핵심 기술 ① 5축 서보 + 변위센서 동기화 (매니저 2순위)
# ============================================================
narration_s7 = """매니저께서 강조하신 첫 번째 핵심 기술, 5축 서보 위치 제어와 아날로그 변위센서 동기화 오차 최소화에 대한 저희의 자산을 말씀드립니다. 직전 양산인 대한전선 WINDER 컨트롤러는 권취기 정밀 장력과 위치 피드백을 RS-485 산업 통신으로 동기화한 사례이며, 5축 동기 제어의 핵심 패턴이 동일합니다. 라즈베리파이 CM4 EtherCAT MFC 컨트롤러는 산업 결정론적 사이클을 양산 환경에서 운영 중입니다. 비접촉 센서 분야에서는 PVDF 압전 기반 층간소음 디바이스와 액체 비접촉 측정 수위 센서로 신호 처리와 노이즈 보정 패턴을 검증했습니다. 변위센서 데이터를 서보 위치와 동기화하는 작업은 측정 사이클 안에서 PLC 인터럽트 우선순위와 EtherCAT 또는 고속 아날로그 입력 모듈의 트리거 동기를 설계하는 문제로, 저희가 양산에서 직접 다뤄온 영역입니다."""

s = prs.slides.add_slide(blank_layout)
add_page_header(s, 6, TOTAL, "핵심 기술 ①")
add_page_footer(s)
add_section_title(s, "5축 서보  ⟷  비접촉 변위센서  동기화",
                  "위치 ↔ 측정값 동기 오차 최소화 (매니저 요구 2순위 ①)")

# 좌측 — 동기화 패턴 다이어그램 (텍스트형 박스 3단계)
left_x = Inches(0.6)
diag_y = Inches(2.4)
diag_w = Inches(6.0)
diag_h = Inches(4.5)
add_card(s, left_x, diag_y, diag_w, diag_h, fill=CREAM)
add_text(s, left_x + Inches(0.3), diag_y + Inches(0.2), diag_w, Inches(0.45),
         "동기화 설계 패턴 — 3-Layer",
         size=16, bold=True, color=NAVY)

# layer 1
ly = diag_y + Inches(0.85)
add_rect(s, left_x + Inches(0.3), ly, diag_w - Inches(0.6), Inches(1.0),
         fill=NAVY)
add_text(s, left_x + Inches(0.45), ly + Inches(0.1), Inches(5.5), Inches(0.4),
         "Layer 1 — EtherCAT / 고속 아날로그 모듈",
         size=13, bold=True, color=WHITE)
add_text(s, left_x + Inches(0.45), ly + Inches(0.45), Inches(5.5), Inches(0.5),
         "5축 서보 위치 = 측정 트리거 = 변위센서 샘플 ── 동일 사이클 캡처",
         size=11, color=CREAM)

# layer 2
ly = diag_y + Inches(2.05)
add_rect(s, left_x + Inches(0.3), ly, diag_w - Inches(0.6), Inches(1.0),
         fill=NAVY_LIGHT)
add_text(s, left_x + Inches(0.45), ly + Inches(0.1), Inches(5.5), Inches(0.4),
         "Layer 2 — PLC 인터럽트 우선순위 / 결정론 사이클",
         size=13, bold=True, color=WHITE)
add_text(s, left_x + Inches(0.45), ly + Inches(0.45), Inches(5.5), Inches(0.5),
         "측정 인터럽트 ＞ 모니터링 ＞ HMI 통신  ── 위치 ↔ 값 오프셋 최소화",
         size=11, color=CREAM)

# layer 3
ly = diag_y + Inches(3.25)
add_rect(s, left_x + Inches(0.3), ly, diag_w - Inches(0.6), Inches(1.0),
         fill=ORANGE)
add_text(s, left_x + Inches(0.45), ly + Inches(0.1), Inches(5.5), Inches(0.4),
         "Layer 3 — 신호 처리 / 노이즈 보정",
         size=13, bold=True, color=WHITE)
add_text(s, left_x + Inches(0.45), ly + Inches(0.45), Inches(5.5), Inches(0.5),
         "이동 평균 + 칼만 필터 + 위치 보간  ── 측정 정밀도 보장",
         size=11, color=CREAM)

# 우측 — 양산 자산 매칭
right_x = Inches(6.95)
right_w = Inches(6.0)
add_card(s, right_x, diag_y, right_w, diag_h, fill=WHITE)
add_rect(s, right_x, diag_y, right_w, Inches(0.5), fill=NAVY)
add_text(s, right_x + Inches(0.2), diag_y + Inches(0.08), right_w, Inches(0.4),
         "직결 양산 자산",
         size=15, bold=True, color=WHITE)

assets = """● 대한전선 WINDER CONTROLLER
   권취기 정밀 장력 + 위치 피드백 양산
   CAN / RS-485 동기 통신 — 5축 패턴 동일

● RPi CM4 EtherCAT MFC 컨트롤러
   산업 결정론 사이클 (100μs급) 양산 운영

● PVDF 압전 층간소음 (위시캣 #154780)
   비접촉 변위 신호 처리 + 노이즈 보정

● 수위센서 QDY30A-B 비접촉 측정
   액체 거리 측정 인접 경험

● STM32F756 컴프레서 밸브 양산
   FreeRTOS 우선순위 + 측정 인터럽트 설계 KC 인증"""
add_text(s, right_x + Inches(0.3), diag_y + Inches(0.7),
         right_w - Inches(0.4), diag_h - Inches(0.8),
         assets, size=12, color=GRAY_TEXT)

add_notes(s, narration_s7)

# ============================================================
# Slide 8 — 핵심 기술 ② 노이즈 차폐 + EMI/EMC
# ============================================================
narration_s8 = """서보 모터 구동 시 발생하는 노이즈가 아날로그 변위센서 데이터에 영향을 주지 않도록 하는 설계는 산업 양산의 가장 중요한 디테일입니다. 저희는 KC 인증을 받은 STM32F756 컴프레서 밸브 컨트롤러에서 산업 EMI·EMC 환경 양산 검증을 거쳤습니다. 표준 대응 4종을 적용합니다. 첫째, 서보 동력 케이블과 센서 신호 케이블의 분리 포설과 쉴드 단일점 접지. 둘째, 아날로그 입력단의 차분 입력과 적정 차단 주파수의 저역 통과 필터. 셋째, PLC와 측정 모듈 사이 광 절연과 신호 그라운드 격리. 넷째, 서보 드라이버 입출력에 페라이트 코어와 노이즈 필터 적용. 일본 TELEC와 유럽 CE 인증을 통과한 3개국 양산 사례가 본 패턴의 신뢰성을 뒷받침합니다."""

s = prs.slides.add_slide(blank_layout)
add_page_header(s, 7, TOTAL, "핵심 기술 ②")
add_page_footer(s)
add_section_title(s, "노이즈 차폐 + EMI / EMC 표준 대응 4종",
                  "서보 노이즈 ✗→ 아날로그 변위센서  (매니저 요구 2순위 ②)")

# 4 quadrant cards
items = [
    ("①  케이블 포설 분리 + 쉴드 단일점 접지",
     "동력선과 신호선 트레이 분리\n쉴드선 단일점 접지 (1-point grounding)\nRS-485 양산 다수에서 검증",
     NAVY),
    ("②  차분 입력 + 저역 통과 필터",
     "아날로그 입력단 차분 (differential) 구성\n적정 차단 주파수 LPF (1~10kHz)\n서보 PWM 고주파 노이즈 차단",
     NAVY_LIGHT),
    ("③  광 절연 + 신호 그라운드 격리",
     "PLC ↔ 측정 모듈 광 커플러 절연\n신호 그라운드 / 동력 그라운드 격리\n그라운드 루프 차단",
     ORANGE),
    ("④  페라이트 코어 + 노이즈 필터",
     "서보 드라이버 입출력 페라이트\n전원 라인 EMI 필터\nKC 인증 양산 표준",
     ORANGE),
]
cell_w = Inches(6.05)
cell_h = Inches(2.05)
grid_x = Inches(0.6)
grid_y = Inches(2.3)
gap_x = Inches(0.18)
gap_y = Inches(0.18)
for i, (title, desc, color) in enumerate(items):
    r = i // 2
    c = i % 2
    x = grid_x + c * (cell_w + gap_x)
    y = grid_y + r * (cell_h + gap_y)
    add_card(s, x, y, cell_w, cell_h, fill=WHITE)
    add_rect(s, x, y, Inches(0.15), cell_h, fill=color)
    add_text(s, x + Inches(0.3), y + Inches(0.15), cell_w - Inches(0.4), Inches(0.5),
             title, size=15, bold=True, color=NAVY)
    add_text(s, x + Inches(0.3), y + Inches(0.7), cell_w - Inches(0.4), Inches(1.2),
             desc, size=12, color=GRAY_TEXT)

# bottom: 3개국 인증 검증
add_text(s, Inches(0.6), Inches(6.7), Inches(12), Inches(0.4),
         "KC · TELEC · CE 3개국 인증 양산  =  본 패턴 신뢰성 검증",
         size=14, bold=True, color=NAVY, align=PP_ALIGN.CENTER)

add_notes(s, narration_s8)

# ============================================================
# Slide 9 — SCADA 5요소 라이브
# ============================================================
narration_s9 = """SCADA 운영 자산을 보여드립니다. 저희는 uttec-sensor.duckdns.org 도메인으로 스마트팩토리 센서 모니터링 시스템을 라이브 운영하고 있습니다. SCADA의 핵심 5요소인 Tag Server, HMI, Recipe, Historian, Alarm을 모두 자체 구현하여 운영 중입니다. Three.js 기반 3D 대시보드 25개 데모와 FastAPI 백엔드, PostgreSQL과 TimescaleDB 시계열 데이터베이스가 결합되어 분체·파쇄·혼합 산업 모니터링에 적용되었으며, 한국기계 공장자동화 협력 사례를 통해 검증되었습니다. 실시간 트렌드 그래프와 측정 완료 데이터의 엑셀 CSV 자동 출력 리포트는 표준 패턴으로 보유하고 있습니다. 클라이언트께서 선호하시는 Cimon, iX Developer, WinCC 같은 상용 SCADA 패키지로의 적응은 OPC UA, Modbus TCP, SQL Historian 같은 산업 표준 인터페이스 기반으로 1주일 내에 완료 가능합니다."""

s = prs.slides.add_slide(blank_layout)
add_page_header(s, 8, TOTAL, "SCADA 운영 자산")
add_page_footer(s)
add_section_title(s, "SCADA 5요소  ·  라이브 운영 중",
                  "uttec-sensor.duckdns.org  ·  실시간 트렌드 + CSV 리포트")

# 5 elements horizontal row
elem_y = Inches(2.4)
elem_h = Inches(1.6)
elem_w = Inches(2.43)
elem_gap = Inches(0.13)
elem_x = Inches(0.6)
elements = [
    ("Tag\nServer", "실시간 값 수집", NAVY),
    ("HMI", "Three.js 3D\n대시보드 25개", NAVY_LIGHT),
    ("Recipe", "제품별 파라미터\n셋 DB", ORANGE),
    ("Historian", "PostgreSQL +\nTimescaleDB", NAVY_LIGHT),
    ("Alarm", "임계값 + n8n /\nGmail 통보", NAVY),
]
for i, (label, desc, color) in enumerate(elements):
    x = elem_x + i * (elem_w + elem_gap)
    add_card(s, x, elem_y, elem_w, elem_h, fill=WHITE)
    add_rect(s, x, elem_y, elem_w, Inches(0.6), fill=color)
    add_text(s, x, elem_y + Inches(0.08), elem_w, Inches(0.5),
             label, size=15, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text(s, x + Inches(0.1), elem_y + Inches(0.72), elem_w - Inches(0.2), Inches(0.8),
             desc, size=12, color=GRAY_TEXT, align=PP_ALIGN.CENTER)

# Lower section — 실시간 트렌드 + CSV 리포트 (left) / 상용 적응 (right)
lower_y = Inches(4.4)
lw = Inches(6.05)
lh = Inches(2.5)

# left card
add_card(s, Inches(0.6), lower_y, lw, lh, fill=CREAM)
add_text(s, Inches(0.85), lower_y + Inches(0.15), lw, Inches(0.5),
         "실시간 트렌드  +  CSV 리포트 자동화",
         size=15, bold=True, color=NAVY)
left_bullets = """●  실시간 두께 데이터 트렌드 그래프 — 표준 컴포넌트
●  측정 완료 후 지정 양식 엑셀(CSV) 자동 출력
●  Historian 기반 품질 트레이스빌리티 보고서
●  한국기계 공장자동화 협력 적용 사례 검증
●  분체 / 파쇄 / 혼합 산업 모니터링 운영"""
add_text(s, Inches(0.85), lower_y + Inches(0.65), lw - Inches(0.4), Inches(1.8),
         left_bullets, size=12, color=GRAY_TEXT)

# right card
rx = Inches(6.95)
add_card(s, rx, lower_y, lw, lh, fill=WHITE)
add_text(s, rx + Inches(0.25), lower_y + Inches(0.15), lw, Inches(0.5),
         "상용 SCADA 1주 적응 가능",
         size=15, bold=True, color=NAVY)
right_bullets = """●  Cimon Ultimate Access / iX Developer / WinCC
●  표준 인터페이스 — OPC UA / Modbus TCP / SQL
●  5요소 자체 구현 자산 → 상용 매핑 1:1
●  XG5000 + Cimon (LS 스택) 즉시 대응
●  선호 플랫폼 확정 시 사전 학습 + 라이선스 산정"""
add_text(s, rx + Inches(0.25), lower_y + Inches(0.65), lw - Inches(0.4), Inches(1.8),
         right_bullets, size=12, color=GRAY_TEXT)

add_notes(s, narration_s9)

# ============================================================
# Slide 11 — TCO 분리 + 3-Plan
# ============================================================
narration_s11 = """매니저께서 가장 중요하게 강조하신 비용 분리와 현실적인 제어기기 역제안 부분입니다. 먼저 RFP에 명시된 대로 2,000만원은 순수 인건비와 개발비 기준으로 픽스합니다. 자재비 BOM과 SCADA 라이선스 비용은 별도로 산출하여 클라이언트께서 총소유비용 TCO를 명확히 인지하실 수 있도록 투명하게 안내드립니다. 제어기기 스택은 세 가지 안을 역제안합니다. Plan A는 클라이언트께서 선호하시는 미쓰비시 PLC와 GT Designer 조합으로 안정성과 호환성을 최우선으로 합니다. Plan B는 LS산전 XGT와 Cimon 조합으로 동등 기능을 15에서 25퍼센트 자재비 절감으로 구현하는 가성비 안입니다. Plan C는 자체 솔루션으로 상용 SCADA 라이선스를 절감하는 옵션으로, 미선정 시에도 부분 적용 가능한 자산입니다. 미팅에서 클라이언트의 우선순위를 확인한 후 최종 단가표를 확정 제출드리겠습니다."""

s = prs.slides.add_slide(blank_layout)
add_page_header(s, 9, TOTAL, "비용 분리 + 3-Plan")
add_page_footer(s)
add_section_title(s, "2,000만원 인건비 Fix  +  3-Plan 역제안",
                  "자재비 TCO 투명 안내 + 현실적 제어기기 역제안 (매니저 요구 3순위)")

# Top fix card
fix_y = Inches(2.25)
add_card(s, Inches(0.6), fix_y, Inches(12.1), Inches(0.75), fill=NAVY)
add_text(s, Inches(0.85), fix_y + Inches(0.1), Inches(12), Inches(0.45),
         "인건비 / 개발비  =  20,000,000원  (Fix)   ·   자재비 (BOM) + SCADA 라이선스  =  별도 (실비)",
         size=15, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
add_text(s, Inches(0.85), fix_y + Inches(0.42), Inches(12), Inches(0.3),
         "RFP 명시 조건과 동일 — 총소유비용(TCO) 명확화로 사후 분쟁 차단",
         size=11, color=CREAM, align=PP_ALIGN.CENTER)

# Three plans
plans_y = Inches(3.25)
plans_h = Inches(3.5)
plan_w = Inches(3.95)
plan_gap = Inches(0.2)

plans = [
    ("Plan A", "미쓰비시 (선호 1순위)",
     "PLC : Mitsubishi FX5 / Q\nSCADA : GT Designer3",
     "● 클라이언트 선호 1순위\n● 안정성 / 호환성 최우선\n● 표준 미쓰비시 생태계",
     "자재비 + 라이선스 별도 산출",
     NAVY, ORANGE),
    ("Plan B", "LS XGT + Cimon (가성비)",
     "PLC : LS XGT (XGB / XGI)\nSCADA : Cimon Ultimate",
     "● 동등 기능 / 국산 스택\n● 자재비 15~25% 절감\n● 사후 유지보수 국내 우수",
     "Plan A 대비 ↓ 15~25%",
     ORANGE, NAVY),
    ("Plan C", "자체 솔루션 (옵션)",
     "PLC : LS XGT / Mitsubishi\nSCADA : UTTEC 자체 (Three.js + DB)",
     "● 상용 SCADA 라이선스 절감\n● 5요소 자체 구현 자산 활용\n● 미선정 시 부분 옵션 가능",
     "Plan A 대비 ↓ 40~55%",
     GREEN_OK, NAVY),
]
for i, (badge, name, stack, bullets, cost, head_color, accent_color) in enumerate(plans):
    x = Inches(0.6) + i * (plan_w + plan_gap)
    add_card(s, x, plans_y, plan_w, plans_h, fill=WHITE)
    add_rect(s, x, plans_y, plan_w, Inches(0.55), fill=head_color)
    add_text(s, x + Inches(0.15), plans_y + Inches(0.08), plan_w - Inches(0.3), Inches(0.4),
             badge, size=13, bold=True, color=WHITE)
    add_text(s, x + Inches(0.15), plans_y + Inches(0.7), plan_w - Inches(0.3), Inches(0.5),
             name, size=15, bold=True, color=NAVY)
    add_text(s, x + Inches(0.15), plans_y + Inches(1.2), plan_w - Inches(0.3), Inches(0.7),
             stack, size=11, color=GRAY_TEXT)
    add_text(s, x + Inches(0.15), plans_y + Inches(2.0), plan_w - Inches(0.3), Inches(1.1),
             bullets, size=11, color=GRAY_TEXT)
    add_chip(s, x + Inches(0.15), plans_y + Inches(2.9), plan_w - Inches(0.3), Inches(0.4),
             cost, accent_color, WHITE, size=11)

add_notes(s, narration_s11)

# ============================================================
# Slide 12 — 90일 일정 + 클로징
# ============================================================
narration_s12 = """90일 일정을 4단계로 제안드립니다. Phase 1은 1일에서 10일, 현장 1~2회 방문으로 1호기 분석과 측정 사양 파악, 플랫폼 선정, 배선도 1차안을 완성합니다. Phase 2는 11일에서 40일, 원격 개발 단계로 PLC 프로그램과 SCADA 화면, Recipe 구조를 개발하고 주 1회 화상 진행 공유와 현장 2~3회 시운전을 진행합니다. Phase 3은 41일에서 70일, 현장 집중 단계로 전장 배선과 통합 시운전, 서보 모션 캘리브레이션을 5~7회 출장으로 완수합니다. Phase 4는 71일에서 90일, 안정화와 Recipe 셋업, 데이터 로깅 검증, 인수인계와 운영자 교육 단계입니다. 당사 위치인 용인 기흥에서 부천 오정구 현장까지 1.5시간으로 출장 협의에 우호적입니다. 양산 6종 제품과 일본 3,800대 수출 실적, KC와 TELEC, CE 3개국 인증 자산으로 안정성, 신뢰성, 비용 투명성 세 가지를 약속드립니다. 미팅 자리에서 더 깊은 협의 기대하겠습니다. 감사합니다."""

s = prs.slides.add_slide(blank_layout)
add_page_header(s, 10, TOTAL, "90일 일정 + 약속")
add_page_footer(s)
add_section_title(s, "90일 4-Phase 일정  +  세 가지 약속",
                  "현장 분석 → 개발 → 시운전 → 안정화 / 인수인계")

# 4 phases horizontal
ph_y = Inches(2.3)
ph_h = Inches(2.2)
ph_w = Inches(3.0)
ph_gap = Inches(0.13)
ph_x = Inches(0.6)
phases = [
    ("Phase 1", "1 ~ 10일", "현장 분석 + 설계 1차",
     "현장 1~2회\n1호기 분석 / 측정 사양\n플랫폼 선정 / 배선도 1차",
     NAVY),
    ("Phase 2", "11 ~ 40일", "PLC + SCADA 1차 개발",
     "원격 개발\nPLC + SCADA + Recipe\n주 1회 화상 + 현장 2~3회",
     NAVY_LIGHT),
    ("Phase 3", "41 ~ 70일", "전장 배선 + 시운전",
     "현장 집중 5~7회\n패널 / 배선 / 캘리브레이션\n1품종 → 다품종 확장",
     ORANGE),
    ("Phase 4", "71 ~ 90일", "안정화 + 인수인계",
     "Alarm 셋업 / 24/7 검증\n운영자 매뉴얼 + 교육\nKC/EMI 사전 점검",
     GREEN_OK),
]
for i, (badge, days, title, desc, color) in enumerate(phases):
    x = ph_x + i * (ph_w + ph_gap)
    add_card(s, x, ph_y, ph_w, ph_h, fill=WHITE)
    add_rect(s, x, ph_y, ph_w, Inches(0.5), fill=color)
    add_text(s, x + Inches(0.15), ph_y + Inches(0.08), ph_w - Inches(0.3), Inches(0.35),
             f"{badge}  ·  {days}", size=12, bold=True, color=WHITE)
    add_text(s, x + Inches(0.15), ph_y + Inches(0.6), ph_w - Inches(0.3), Inches(0.5),
             title, size=14, bold=True, color=NAVY)
    add_text(s, x + Inches(0.15), ph_y + Inches(1.1), ph_w - Inches(0.3), Inches(1.0),
             desc, size=11, color=GRAY_TEXT)

# 3 promises
pr_y = Inches(4.85)
pr_h = Inches(0.95)
pr_w = Inches(3.95)
pr_gap = Inches(0.2)
promises = [
    ("안정성", "양산 6종 + 일본 3,800대 검증", NAVY),
    ("신뢰성", "KC + TELEC + CE 3개국 인증", ORANGE),
    ("비용 투명성", "인건비 Fix + 자재비 별도 + 3-Plan", GREEN_OK),
]
for i, (head, body, color) in enumerate(promises):
    x = Inches(0.6) + i * (pr_w + pr_gap)
    add_card(s, x, pr_y, pr_w, pr_h, fill=WHITE)
    add_rect(s, x, pr_y, Inches(0.18), pr_h, fill=color)
    add_text(s, x + Inches(0.35), pr_y + Inches(0.1), pr_w - Inches(0.4), Inches(0.4),
             head, size=15, bold=True, color=color)
    add_text(s, x + Inches(0.35), pr_y + Inches(0.5), pr_w - Inches(0.4), Inches(0.45),
             body, size=11, color=GRAY_TEXT)

# Closing
cl_y = Inches(6.05)
add_card(s, Inches(0.6), cl_y, Inches(12.1), Inches(1.0), fill=NAVY)
add_text(s, Inches(0.6), cl_y + Inches(0.12), Inches(12.1), Inches(0.4),
         "용인 기흥 → 부천 오정구  1.5시간  ·  출장 협의 우호적",
         size=14, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
add_text(s, Inches(0.6), cl_y + Inches(0.55), Inches(12.1), Inches(0.4),
         "미팅 자리에서 더 깊은 협의 기대하겠습니다.  감사합니다.",
         size=16, bold=True, color=ORANGE, align=PP_ALIGN.CENTER)

add_notes(s, narration_s12)

# ============================================================
# Save
# ============================================================
out = r"C:\todo\today\동영상\동아정밀_미팅\동아정밀_PET두께측정기_제안.pptx"
prs.save(out)
print(f"Saved: {out}")
print(f"Slides: {len(prs.slides)}")

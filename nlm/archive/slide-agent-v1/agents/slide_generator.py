"""
Slide Generator - PPTX 슬라이드 생성기
3D 렌더링 스타일의 시네마틱 디자인
"""

import json
from datetime import datetime
from pathlib import Path
from typing import Optional
import logging

logger = logging.getLogger(__name__)

# Lazy import for pptx
Presentation = None
Inches = None
Pt = None
Emu = None
RgbColor = None
PP_ALIGN = None
MSO_ANCHOR = None
MSO_SHAPE = None
PPTX_AVAILABLE = False


def _load_pptx():
    """Lazy load python-pptx"""
    global Presentation, Inches, Pt, Emu, RgbColor, PP_ALIGN, MSO_ANCHOR, MSO_SHAPE, PPTX_AVAILABLE
    if PPTX_AVAILABLE:
        return True
    try:
        from pptx import Presentation as _Presentation
        from pptx.util import Inches as _Inches, Pt as _Pt, Emu as _Emu
        from pptx.dml.color import RGBColor as _RgbColor
        from pptx.enum.text import PP_ALIGN as _PP_ALIGN, MSO_ANCHOR as _MSO_ANCHOR
        from pptx.enum.shapes import MSO_SHAPE as _MSO_SHAPE
        Presentation = _Presentation
        Inches = _Inches
        Pt = _Pt
        Emu = _Emu
        RgbColor = _RgbColor
        PP_ALIGN = _PP_ALIGN
        MSO_ANCHOR = _MSO_ANCHOR
        MSO_SHAPE = _MSO_SHAPE
        PPTX_AVAILABLE = True
        return True
    except ImportError as e:
        logger.error(f"python-pptx not installed or import error: {e}")
        return False


class SlideGenerator:
    """
    PPTX 슬라이드 생성기

    디자인 스타일: 3D 렌더링 / 시네마틱
    - 배경: #1A1A2E (다크 네이비)
    - 텍스트: #FFFFFF (화이트)
    - 강조: #6366F1 (인디고)
    - 폰트: Sora (없으면 Arial로 대체)
    """

    def __init__(self, output_dir: str):
        self.output_dir = Path(output_dir)
        self.slides_dir = self.output_dir / "slides"
        self.slides_dir.mkdir(parents=True, exist_ok=True)
        self.COLORS = {}
        self.FONT_NAME = "Sora"  # 3D/크리에이티브 스타일 폰트

    def _init_colors(self):
        """Initialize 3D rendering style colors"""
        self.COLORS = {
            'background': RgbColor(0x1A, 0x1A, 0x2E),      # #1A1A2E - 다크 네이비
            'text': RgbColor(0xFF, 0xFF, 0xFF),             # #FFFFFF - 화이트
            'accent': RgbColor(0x63, 0x66, 0xF1),           # #6366F1 - 인디고
            'accent_light': RgbColor(0x81, 0x84, 0xF4),     # 밝은 인디고
            'accent_glow': RgbColor(0xA5, 0xB4, 0xFC),      # 글로우 효과
            'text_muted': RgbColor(0x94, 0xA3, 0xB8),       # 뮤트 텍스트
            'gradient_start': RgbColor(0x1A, 0x1A, 0x2E),   # 그라데이션 시작
            'gradient_end': RgbColor(0x0F, 0x0F, 0x1A),     # 그라데이션 끝
            'card_bg': RgbColor(0x25, 0x25, 0x3A),          # 카드 배경
            'border': RgbColor(0x3D, 0x3D, 0x5C),           # 테두리
        }

    def _set_slide_background(self, slide, prs):
        """슬라이드 배경을 다크 네이비로 설정"""
        background = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(0), Inches(0),
            prs.slide_width, prs.slide_height
        )
        background.fill.solid()
        background.fill.fore_color.rgb = self.COLORS['background']
        background.line.fill.background()

        # 배경을 맨 뒤로 보내기
        spTree = slide.shapes._spTree
        sp = background._element
        spTree.remove(sp)
        spTree.insert(2, sp)

    def _add_accent_glow(self, slide, x, y, width, height):
        """강조 글로우 효과 추가"""
        glow = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(x), Inches(y),
            Inches(width), Inches(height)
        )
        glow.fill.solid()
        glow.fill.fore_color.rgb = self.COLORS['accent']
        glow.line.fill.background()
        # 투명도 효과 (python-pptx 제한으로 완전한 투명도는 어려움)
        return glow

    def _add_card(self, slide, x, y, width, height):
        """3D 스타일 카드 추가"""
        card = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(x), Inches(y),
            Inches(width), Inches(height)
        )
        card.fill.solid()
        card.fill.fore_color.rgb = self.COLORS['card_bg']
        card.line.color.rgb = self.COLORS['border']
        card.line.width = Pt(1)
        return card

    def create_presentation(self, slide_content: dict) -> Optional[str]:
        """
        3D 렌더링 스타일 프레젠테이션 생성
        """
        if not _load_pptx():
            logger.error("python-pptx not available")
            return None

        self._init_colors()

        prs = Presentation()
        prs.slide_width = Inches(13.333)  # 16:9
        prs.slide_height = Inches(7.5)

        # 1. 타이틀 슬라이드
        self._add_title_slide(prs, slide_content)

        # 2. 목차 슬라이드
        self._add_agenda_slide(prs, slide_content)

        # 3. 소스 개요 슬라이드
        self._add_sources_slide(prs, slide_content)

        # 4. 각 섹션별 슬라이드
        for section in slide_content.get('sections', []):
            self._add_section_slide(prs, section)

        # 5. 핵심 인사이트 슬라이드
        if slide_content.get('key_insights'):
            self._add_insights_slide(prs, slide_content)

        # 6. 액션 아이템 슬라이드
        if slide_content.get('action_items'):
            self._add_action_items_slide(prs, slide_content)

        # 7. 마무리 슬라이드
        self._add_closing_slide(prs, slide_content)

        # 저장
        date_str = slide_content.get('date', datetime.now().strftime('%Y-%m-%d'))
        output_path = self.slides_dir / f"ai_digest_{date_str.replace('-', '')}_3d.pptx"

        prs.save(str(output_path))
        logger.info(f"Created presentation: {output_path}")

        return str(output_path)

    def _add_title_slide(self, prs, content: dict):
        """시네마틱 타이틀 슬라이드"""
        slide_layout = prs.slide_layouts[6]  # Blank
        slide = prs.slides.add_slide(slide_layout)

        # 다크 배경
        self._set_slide_background(slide, prs)

        # 상단 글로우 액센트 바
        accent_bar = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(0), Inches(0),
            prs.slide_width, Inches(0.1)
        )
        accent_bar.fill.solid()
        accent_bar.fill.fore_color.rgb = self.COLORS['accent']
        accent_bar.line.fill.background()

        # 중앙 글로우 원 (3D 오브젝트 대체)
        glow_circle = slide.shapes.add_shape(
            MSO_SHAPE.OVAL,
            Inches(5.5), Inches(1.5),
            Inches(2.333), Inches(2.333)
        )
        glow_circle.fill.solid()
        glow_circle.fill.fore_color.rgb = self.COLORS['accent']
        glow_circle.line.fill.background()

        # 타이틀
        title_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(4),
            Inches(12.333), Inches(1.2)
        )
        title_frame = title_box.text_frame
        title_para = title_frame.paragraphs[0]
        title_para.text = content.get('title', 'AI News Digest')
        title_para.font.size = Pt(54)
        title_para.font.bold = True
        title_para.font.name = self.FONT_NAME
        title_para.font.color.rgb = self.COLORS['text']
        title_para.alignment = PP_ALIGN.CENTER

        # 날짜
        date_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(5.3),
            Inches(12.333), Inches(0.5)
        )
        date_frame = date_box.text_frame
        date_para = date_frame.paragraphs[0]
        date_para.text = content.get('date', datetime.now().strftime('%Y-%m-%d'))
        date_para.font.size = Pt(20)
        date_para.font.name = self.FONT_NAME
        date_para.font.color.rgb = self.COLORS['accent_glow']
        date_para.alignment = PP_ALIGN.CENTER

        # 소스 수
        sources_count = len(content.get('sources', []))
        subtitle_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(5.9),
            Inches(12.333), Inches(0.5)
        )
        subtitle_frame = subtitle_box.text_frame
        subtitle_para = subtitle_frame.paragraphs[0]
        subtitle_para.text = f"{sources_count} Videos Analyzed"
        subtitle_para.font.size = Pt(16)
        subtitle_para.font.name = self.FONT_NAME
        subtitle_para.font.color.rgb = self.COLORS['text_muted']
        subtitle_para.alignment = PP_ALIGN.CENTER

    def _add_agenda_slide(self, prs, content: dict):
        """목차 슬라이드"""
        slide_layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(slide_layout)

        self._set_slide_background(slide, prs)
        self._add_slide_title(slide, "AGENDA")

        # 목차 아이템
        agenda_items = ["Video Sources"]

        for section in content.get('sections', [])[:4]:
            agenda_items.append(section.get('title', 'Section'))

        if content.get('key_insights'):
            agenda_items.append("Key Insights")

        if content.get('action_items'):
            agenda_items.append("Action Items")

        # 목차 카드들
        y_pos = 1.8
        for i, item in enumerate(agenda_items, 1):
            # 카드
            card = self._add_card(slide, 1, y_pos, 11.333, 0.65)

            # 번호 (강조색)
            num_box = slide.shapes.add_textbox(
                Inches(1.3), Inches(y_pos + 0.12),
                Inches(0.5), Inches(0.4)
            )
            num_frame = num_box.text_frame
            num_para = num_frame.paragraphs[0]
            num_para.text = f"{i:02d}"
            num_para.font.size = Pt(18)
            num_para.font.bold = True
            num_para.font.name = self.FONT_NAME
            num_para.font.color.rgb = self.COLORS['accent']

            # 텍스트
            item_box = slide.shapes.add_textbox(
                Inches(2.2), Inches(y_pos + 0.12),
                Inches(9), Inches(0.4)
            )
            item_frame = item_box.text_frame
            item_para = item_frame.paragraphs[0]
            item_para.text = item
            item_para.font.size = Pt(18)
            item_para.font.name = self.FONT_NAME
            item_para.font.color.rgb = self.COLORS['text']

            y_pos += 0.85

    def _add_sources_slide(self, prs, content: dict):
        """소스 개요 슬라이드"""
        slide_layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(slide_layout)

        self._set_slide_background(slide, prs)
        self._add_slide_title(slide, "VIDEO SOURCES")

        sources = content.get('sources', [])

        # 2열 레이아웃
        col1_x, col2_x = 0.8, 7
        y_positions = [1.8, 3.2, 4.6, 6.0]

        for i, source in enumerate(sources[:8]):
            col_x = col1_x if i < 4 else col2_x
            y_pos = y_positions[i % 4]

            # 카드
            card = self._add_card(slide, col_x, y_pos, 5.5, 1.2)

            # 채널명 (강조색)
            channel_box = slide.shapes.add_textbox(
                Inches(col_x + 0.2), Inches(y_pos + 0.15),
                Inches(5), Inches(0.35)
            )
            channel_frame = channel_box.text_frame
            channel_para = channel_frame.paragraphs[0]
            channel_para.text = source.get('channel', 'Unknown')
            channel_para.font.size = Pt(12)
            channel_para.font.bold = True
            channel_para.font.name = self.FONT_NAME
            channel_para.font.color.rgb = self.COLORS['accent']

            # 영상 제목
            title_box = slide.shapes.add_textbox(
                Inches(col_x + 0.2), Inches(y_pos + 0.55),
                Inches(5), Inches(0.5)
            )
            title_frame = title_box.text_frame
            title_para = title_frame.paragraphs[0]
            title_text = source.get('title', 'Untitled')
            title_para.text = title_text[:45] + ('...' if len(title_text) > 45 else '')
            title_para.font.size = Pt(13)
            title_para.font.name = self.FONT_NAME
            title_para.font.color.rgb = self.COLORS['text']

    def _add_section_slide(self, prs, section: dict):
        """섹션 슬라이드"""
        slide_layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(slide_layout)

        self._set_slide_background(slide, prs)
        self._add_slide_title(slide, section.get('title', 'SECTION').upper())

        content = section.get('content', '')
        section_type = section.get('type', 'general')

        if section_type == 'channel_summary' and section.get('videos'):
            # 채널별 영상 카드
            y_pos = 1.8
            for video in section.get('videos', [])[:4]:
                card = self._add_card(slide, 0.8, y_pos, 11.733, 1.2)

                # 제목
                video_title_box = slide.shapes.add_textbox(
                    Inches(1.1), Inches(y_pos + 0.2),
                    Inches(11), Inches(0.4)
                )
                video_title_frame = video_title_box.text_frame
                video_title_para = video_title_frame.paragraphs[0]
                video_title_para.text = video.get('title', 'Untitled')[:70]
                video_title_para.font.size = Pt(16)
                video_title_para.font.bold = True
                video_title_para.font.name = self.FONT_NAME
                video_title_para.font.color.rgb = self.COLORS['text']

                # 미리보기
                if video.get('transcript_preview'):
                    preview_box = slide.shapes.add_textbox(
                        Inches(1.1), Inches(y_pos + 0.65),
                        Inches(11), Inches(0.4)
                    )
                    preview_frame = preview_box.text_frame
                    preview_para = preview_frame.paragraphs[0]
                    preview_para.text = video.get('transcript_preview', '')[:120] + '...'
                    preview_para.font.size = Pt(11)
                    preview_para.font.name = self.FONT_NAME
                    preview_para.font.color.rgb = self.COLORS['text_muted']

                y_pos += 1.4
        else:
            # 일반 콘텐츠 카드
            card = self._add_card(slide, 0.8, 1.8, 11.733, 5)

            content_box = slide.shapes.add_textbox(
                Inches(1.2), Inches(2.1),
                Inches(11), Inches(4.5)
            )
            content_frame = content_box.text_frame
            content_frame.word_wrap = True

            lines = str(content).split('\n')
            for i, line in enumerate(lines[:12]):
                if i == 0:
                    para = content_frame.paragraphs[0]
                else:
                    para = content_frame.add_paragraph()

                para.text = line.strip()
                para.font.size = Pt(14)
                para.font.name = self.FONT_NAME
                para.font.color.rgb = self.COLORS['text']

    def _add_insights_slide(self, prs, content: dict):
        """핵심 인사이트 슬라이드"""
        slide_layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(slide_layout)

        self._set_slide_background(slide, prs)
        self._add_slide_title(slide, "KEY INSIGHTS")

        insights = content.get('key_insights', [])

        # 메인 카드
        card = self._add_card(slide, 0.8, 1.8, 11.733, 5)

        # 강조 아이콘 (원)
        icon = slide.shapes.add_shape(
            MSO_SHAPE.OVAL,
            Inches(1.2), Inches(2.2),
            Inches(0.5), Inches(0.5)
        )
        icon.fill.solid()
        icon.fill.fore_color.rgb = self.COLORS['accent']
        icon.line.fill.background()

        for insight in insights:
            insight_text = str(insight.get('content', insight) if isinstance(insight, dict) else insight)

            insight_box = slide.shapes.add_textbox(
                Inches(2), Inches(2.2),
                Inches(10), Inches(4.3)
            )
            insight_frame = insight_box.text_frame
            insight_frame.word_wrap = True

            lines = insight_text.split('\n')
            for i, line in enumerate(lines[:10]):
                if i == 0:
                    para = insight_frame.paragraphs[0]
                else:
                    para = insight_frame.add_paragraph()

                para.text = line.strip()
                para.font.size = Pt(15)
                para.font.name = self.FONT_NAME
                para.font.color.rgb = self.COLORS['text']

            break

    def _add_action_items_slide(self, prs, content: dict):
        """액션 아이템 슬라이드"""
        slide_layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(slide_layout)

        self._set_slide_background(slide, prs)
        self._add_slide_title(slide, "ACTION ITEMS")

        action_items = content.get('action_items', [])

        card = self._add_card(slide, 0.8, 1.8, 11.733, 5)

        for item in action_items:
            item_text = str(item.get('content', item) if isinstance(item, dict) else item)

            item_box = slide.shapes.add_textbox(
                Inches(1.2), Inches(2.1),
                Inches(11), Inches(4.5)
            )
            item_frame = item_box.text_frame
            item_frame.word_wrap = True

            lines = item_text.split('\n')
            for i, line in enumerate(lines[:10]):
                if i == 0:
                    para = item_frame.paragraphs[0]
                else:
                    para = item_frame.add_paragraph()

                # 불릿 스타일
                if line.strip().startswith('-'):
                    para.text = "  " + line.strip()
                else:
                    para.text = line.strip()
                para.font.size = Pt(14)
                para.font.name = self.FONT_NAME
                para.font.color.rgb = self.COLORS['text']

            break

    def _add_closing_slide(self, prs, content: dict):
        """마무리 슬라이드 - 시네마틱"""
        slide_layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(slide_layout)

        self._set_slide_background(slide, prs)

        # 중앙 글로우 효과
        glow = slide.shapes.add_shape(
            MSO_SHAPE.OVAL,
            Inches(4.5), Inches(2),
            Inches(4.333), Inches(2.5)
        )
        glow.fill.solid()
        glow.fill.fore_color.rgb = self.COLORS['accent']
        glow.line.fill.background()

        # Thank You 텍스트
        thanks_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(4.8),
            Inches(12.333), Inches(1)
        )
        thanks_frame = thanks_box.text_frame
        thanks_para = thanks_frame.paragraphs[0]
        thanks_para.text = "THANK YOU"
        thanks_para.font.size = Pt(56)
        thanks_para.font.bold = True
        thanks_para.font.name = self.FONT_NAME
        thanks_para.font.color.rgb = self.COLORS['text']
        thanks_para.alignment = PP_ALIGN.CENTER

        # 생성 정보
        info_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(6.2),
            Inches(12.333), Inches(0.5)
        )
        info_frame = info_box.text_frame
        info_para = info_frame.paragraphs[0]
        info_para.text = f"Generated {datetime.now().strftime('%Y-%m-%d %H:%M')}"
        info_para.font.size = Pt(14)
        info_para.font.name = self.FONT_NAME
        info_para.font.color.rgb = self.COLORS['text_muted']
        info_para.alignment = PP_ALIGN.CENTER

        # 하단 액센트 바
        accent_bar = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(0), Inches(7.4),
            prs.slide_width, Inches(0.1)
        )
        accent_bar.fill.solid()
        accent_bar.fill.fore_color.rgb = self.COLORS['accent']
        accent_bar.line.fill.background()

    def _add_slide_title(self, slide, title: str):
        """슬라이드 제목 - 미니멀 스타일"""
        # 강조색 액센트 바
        accent_line = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(0.8), Inches(0.5),
            Inches(0.15), Inches(0.6)
        )
        accent_line.fill.solid()
        accent_line.fill.fore_color.rgb = self.COLORS['accent']
        accent_line.line.fill.background()

        # 제목
        title_box = slide.shapes.add_textbox(
            Inches(1.2), Inches(0.45),
            Inches(11), Inches(0.8)
        )
        title_frame = title_box.text_frame
        title_para = title_frame.paragraphs[0]
        title_para.text = title
        title_para.font.size = Pt(28)
        title_para.font.bold = True
        title_para.font.name = self.FONT_NAME
        title_para.font.color.rgb = self.COLORS['text']


if __name__ == "__main__":
    logging.basicConfig(level=logging.INFO)

    # 테스트 데이터
    test_content = {
        'title': 'AI News Digest',
        'date': '2026-02-18',
        'sources': [
            {'title': 'Test Video 1 - How to Build AI Agents', 'channel': 'Nate Herk', 'url': 'https://youtube.com/1'},
            {'title': 'Test Video 2 - Claude Code Tutorial', 'channel': 'Nick Saraev', 'url': 'https://youtube.com/2'},
        ],
        'sections': [
            {
                'title': 'Key Topics',
                'content': 'This is the key topics section.\n- AI Automation Trends\n- Claude Code Features\n- Workflow Optimization',
                'type': 'topics'
            }
        ],
        'key_insights': ['AI development is accelerating rapidly.\nKey tools include Claude, GPT, and various automation platforms.'],
        'action_items': ['- Start learning Claude Code\n- Build your first AI workflow\n- Join AI communities']
    }

    base_dir = Path(__file__).parent.parent
    generator = SlideGenerator(output_dir=base_dir / "output")

    result = generator.create_presentation(test_content)
    if result:
        print(f"Created: {result}")
    else:
        print("Failed to create presentation")

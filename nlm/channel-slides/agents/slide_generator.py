"""
Slide Generator - 채널별 PPTX 슬라이드 생성기 (최대 10장)
3D 렌더링 스타일의 시네마틱 디자인
"""

import json
from datetime import datetime
from pathlib import Path
from typing import Optional
import logging

logger = logging.getLogger(__name__)

# Lazy import
Presentation = None
Inches = None
Pt = None
RgbColor = None
PP_ALIGN = None
MSO_ANCHOR = None
MSO_SHAPE = None
PPTX_AVAILABLE = False


def _load_pptx():
    """Lazy load python-pptx"""
    global Presentation, Inches, Pt, RgbColor, PP_ALIGN, MSO_ANCHOR, MSO_SHAPE, PPTX_AVAILABLE
    if PPTX_AVAILABLE:
        return True
    try:
        from pptx import Presentation as _Presentation
        from pptx.util import Inches as _Inches, Pt as _Pt
        from pptx.dml.color import RGBColor as _RgbColor
        from pptx.enum.text import PP_ALIGN as _PP_ALIGN, MSO_ANCHOR as _MSO_ANCHOR
        from pptx.enum.shapes import MSO_SHAPE as _MSO_SHAPE
        Presentation = _Presentation
        Inches = _Inches
        Pt = _Pt
        RgbColor = _RgbColor
        PP_ALIGN = _PP_ALIGN
        MSO_ANCHOR = _MSO_ANCHOR
        MSO_SHAPE = _MSO_SHAPE
        PPTX_AVAILABLE = True
        return True
    except ImportError as e:
        logger.error(f"python-pptx import error: {e}")
        return False


class SlideGenerator:
    """
    채널별 PPTX 슬라이드 생성기 (최대 10장)

    디자인: 3D 렌더링 / 시네마틱
    - 배경: #1A1A2E
    - 텍스트: #FFFFFF
    - 강조: #6366F1
    """

    def __init__(self, output_dir: str):
        self.output_dir = Path(output_dir)
        self.slides_dir = self.output_dir / "slides"
        self.slides_dir.mkdir(parents=True, exist_ok=True)
        self.COLORS = {}
        self.FONT_NAME = "Sora"

    def _init_colors(self):
        """컬러 초기화"""
        self.COLORS = {
            'background': RgbColor(0x1A, 0x1A, 0x2E),
            'text': RgbColor(0xFF, 0xFF, 0xFF),
            'accent': RgbColor(0x63, 0x66, 0xF1),
            'accent_glow': RgbColor(0xA5, 0xB4, 0xFC),
            'text_muted': RgbColor(0x94, 0xA3, 0xB8),
            'card_bg': RgbColor(0x25, 0x25, 0x3A),
            'border': RgbColor(0x3D, 0x3D, 0x5C),
        }

    def _set_background(self, slide, prs):
        """다크 배경 설정"""
        bg = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(0), Inches(0),
            prs.slide_width, prs.slide_height
        )
        bg.fill.solid()
        bg.fill.fore_color.rgb = self.COLORS['background']
        bg.line.fill.background()

        spTree = slide.shapes._spTree
        sp = bg._element
        spTree.remove(sp)
        spTree.insert(2, sp)

    def _add_card(self, slide, x, y, width, height):
        """카드 추가"""
        card = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(x), Inches(y),
            Inches(width), Inches(height)
        )
        card.fill.solid()
        card.fill.fore_color.rgb = self.COLORS['card_bg']
        card.line.color.rgb = self.COLORS['border']
        card.line.width = Pt(1)
        return card

    def create_channel_presentation(self, slide_content: dict) -> Optional[str]:
        """
        채널별 PPTX 생성 (최대 10장)
        """
        if not _load_pptx():
            logger.error("python-pptx not available")
            return None

        self._init_colors()

        prs = Presentation()
        prs.slide_width = Inches(13.333)
        prs.slide_height = Inches(7.5)

        slides_data = slide_content.get('slides', [])[:10]  # 최대 10장

        for slide_data in slides_data:
            slide_type = slide_data.get('type', 'content')

            if slide_type == 'title':
                self._add_title_slide(prs, slide_data)
            elif slide_type == 'video_list':
                self._add_video_list_slide(prs, slide_data)
            elif slide_type == 'closing':
                self._add_closing_slide(prs, slide_data)
            else:
                self._add_content_slide(prs, slide_data)

        # 저장
        channel_name = slide_content.get('channel', 'Unknown')
        safe_name = channel_name.replace(' ', '_').replace('/', '_')
        date_str = slide_content.get('date', datetime.now().strftime('%Y-%m-%d'))
        output_path = self.slides_dir / f"{safe_name}_{date_str.replace('-', '')}.pptx"

        prs.save(str(output_path))
        logger.info(f"Created: {output_path.name} ({len(slides_data)} slides)")

        return str(output_path)

    def _add_title_slide(self, prs, data: dict):
        """타이틀 슬라이드"""
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        self._set_background(slide, prs)

        # 상단 액센트 바
        bar = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(0), Inches(0),
            prs.slide_width, Inches(0.1)
        )
        bar.fill.solid()
        bar.fill.fore_color.rgb = self.COLORS['accent']
        bar.line.fill.background()

        # 글로우 원
        circle = slide.shapes.add_shape(
            MSO_SHAPE.OVAL,
            Inches(5.5), Inches(1.2),
            Inches(2.333), Inches(2.333)
        )
        circle.fill.solid()
        circle.fill.fore_color.rgb = self.COLORS['accent']
        circle.line.fill.background()

        # 타이틀
        title_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(4),
            Inches(12.333), Inches(1.2)
        )
        p = title_box.text_frame.paragraphs[0]
        p.text = data.get('title', 'Channel Insights')
        p.font.size = Pt(48)
        p.font.bold = True
        p.font.name = self.FONT_NAME
        p.font.color.rgb = self.COLORS['text']
        p.alignment = PP_ALIGN.CENTER

        # 서브타이틀
        sub_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(5.3),
            Inches(12.333), Inches(0.5)
        )
        p = sub_box.text_frame.paragraphs[0]
        p.text = data.get('subtitle', '')
        p.font.size = Pt(20)
        p.font.name = self.FONT_NAME
        p.font.color.rgb = self.COLORS['accent_glow']
        p.alignment = PP_ALIGN.CENTER

        # 날짜
        date_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(6),
            Inches(12.333), Inches(0.4)
        )
        p = date_box.text_frame.paragraphs[0]
        p.text = data.get('date', '')
        p.font.size = Pt(14)
        p.font.name = self.FONT_NAME
        p.font.color.rgb = self.COLORS['text_muted']
        p.alignment = PP_ALIGN.CENTER

    def _add_video_list_slide(self, prs, data: dict):
        """영상 목록 슬라이드"""
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        self._set_background(slide, prs)
        self._add_slide_title(slide, data.get('title', 'Videos'))

        items = data.get('items', [])[:8]  # 최대 8개
        y_pos = 1.8

        for i, item in enumerate(items, 1):
            card = self._add_card(slide, 1, y_pos, 11.333, 0.6)

            # 번호
            num_box = slide.shapes.add_textbox(
                Inches(1.3), Inches(y_pos + 0.1),
                Inches(0.5), Inches(0.4)
            )
            p = num_box.text_frame.paragraphs[0]
            p.text = f"{i:02d}"
            p.font.size = Pt(16)
            p.font.bold = True
            p.font.name = self.FONT_NAME
            p.font.color.rgb = self.COLORS['accent']

            # 제목
            title_box = slide.shapes.add_textbox(
                Inches(2.1), Inches(y_pos + 0.1),
                Inches(9.5), Inches(0.4)
            )
            p = title_box.text_frame.paragraphs[0]
            p.text = item[:70] if len(item) > 70 else item
            p.font.size = Pt(14)
            p.font.name = self.FONT_NAME
            p.font.color.rgb = self.COLORS['text']

            y_pos += 0.7

    def _add_content_slide(self, prs, data: dict):
        """콘텐츠 슬라이드"""
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        self._set_background(slide, prs)
        self._add_slide_title(slide, data.get('title', 'Content'))

        # 카드
        card = self._add_card(slide, 0.8, 1.6, 11.733, 5.4)

        # 콘텐츠
        content = data.get('content', '')
        content_box = slide.shapes.add_textbox(
            Inches(1.2), Inches(1.9),
            Inches(11), Inches(5)
        )
        tf = content_box.text_frame
        tf.word_wrap = True

        # 줄 단위로 처리
        lines = str(content).split('\n')
        for i, line in enumerate(lines[:18]):  # 최대 18줄
            if i == 0:
                para = tf.paragraphs[0]
            else:
                para = tf.add_paragraph()

            para.text = line.strip()
            para.font.size = Pt(13)
            para.font.name = self.FONT_NAME
            para.font.color.rgb = self.COLORS['text']

    def _add_closing_slide(self, prs, data: dict):
        """마무리 슬라이드"""
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        self._set_background(slide, prs)

        # 글로우
        glow = slide.shapes.add_shape(
            MSO_SHAPE.OVAL,
            Inches(4.5), Inches(2),
            Inches(4.333), Inches(2.5)
        )
        glow.fill.solid()
        glow.fill.fore_color.rgb = self.COLORS['accent']
        glow.line.fill.background()

        # Thank You
        thanks_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(4.8),
            Inches(12.333), Inches(1)
        )
        p = thanks_box.text_frame.paragraphs[0]
        p.text = data.get('title', 'THANK YOU')
        p.font.size = Pt(48)
        p.font.bold = True
        p.font.name = self.FONT_NAME
        p.font.color.rgb = self.COLORS['text']
        p.alignment = PP_ALIGN.CENTER

        # 서브타이틀
        sub_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(6),
            Inches(12.333), Inches(0.4)
        )
        p = sub_box.text_frame.paragraphs[0]
        p.text = data.get('subtitle', '')
        p.font.size = Pt(14)
        p.font.name = self.FONT_NAME
        p.font.color.rgb = self.COLORS['text_muted']
        p.alignment = PP_ALIGN.CENTER

        # 하단 바
        bar = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(0), Inches(7.4),
            prs.slide_width, Inches(0.1)
        )
        bar.fill.solid()
        bar.fill.fore_color.rgb = self.COLORS['accent']
        bar.line.fill.background()

    def _add_slide_title(self, slide, title: str):
        """슬라이드 제목"""
        # 액센트 바
        bar = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(0.8), Inches(0.5),
            Inches(0.15), Inches(0.55)
        )
        bar.fill.solid()
        bar.fill.fore_color.rgb = self.COLORS['accent']
        bar.line.fill.background()

        # 제목
        title_box = slide.shapes.add_textbox(
            Inches(1.2), Inches(0.45),
            Inches(11), Inches(0.7)
        )
        p = title_box.text_frame.paragraphs[0]
        p.text = title.upper()
        p.font.size = Pt(26)
        p.font.bold = True
        p.font.name = self.FONT_NAME
        p.font.color.rgb = self.COLORS['text']

    def create_all_presentations(self, channel_contents: dict[str, dict]) -> dict[str, str]:
        """
        모든 채널의 프레젠테이션 생성

        Returns:
            {channel_name: pptx_path}
        """
        results = {}

        for channel_name, slide_content in channel_contents.items():
            logger.info(f"Generating slides for: {channel_name}")
            pptx_path = self.create_channel_presentation(slide_content)

            if pptx_path:
                results[channel_name] = pptx_path

        logger.info(f"\nGenerated {len(results)} presentations")
        return results


if __name__ == "__main__":
    logging.basicConfig(level=logging.INFO, format='%(message)s')

    # 테스트
    test_content = {
        'channel': 'Test Channel',
        'title': 'Test Channel',
        'date': '2026-02-18',
        'video_count': 2,
        'slides': [
            {'type': 'title', 'title': 'Test Channel', 'subtitle': '2 Videos', 'date': '2026-02-18'},
            {'type': 'video_list', 'title': 'Videos', 'items': ['Video 1', 'Video 2']},
            {'type': 'content', 'title': 'Main Topics', 'content': 'Topic 1\nTopic 2\nTopic 3'},
            {'type': 'closing', 'title': 'Thank You', 'subtitle': 'Analysis Complete'}
        ]
    }

    base_dir = Path(__file__).parent.parent
    generator = SlideGenerator(output_dir=base_dir / "output")

    result = generator.create_channel_presentation(test_content)
    print(f"Created: {result}")

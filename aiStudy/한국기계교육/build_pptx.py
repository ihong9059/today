"""한국기계 AI 교육 소개 PPT (10장) 생성
- 1920x1080 (16:9 widescreen)
- 한국어 / 비즈니스 톤
- 각 슬라이드 notes에 narration
- 동아정밀 build_pptx.py 패턴 재사용 (UTTEC 위시캣 미팅 자료 SOP 2번째 적용 사례)
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

# ============================================================
# Brand palette
# ============================================================
NAVY = RGBColor(0x10, 0x2A, 0x4C)
NAVY_LIGHT = RGBColor(0x1E, 0x3A, 0x5F)
ORANGE = RGBColor(0xF1, 0x7A, 0x21)
GREEN_OK = RGBColor(0x2E, 0x8B, 0x57)
RED_FOCUS = RGBColor(0xC0, 0x39, 0x2B)
CREAM = RGBColor(0xFA, 0xF6, 0xEC)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
GRAY_TEXT = RGBColor(0x33, 0x33, 0x33)
GRAY_LIGHT = RGBColor(0xE0, 0xE0, 0xE0)
GRAY_MID = RGBColor(0x88, 0x88, 0x88)
GREEN_LIGHT = RGBColor(0xE8, 0xF5, 0xEC)
ORANGE_LIGHT = RGBColor(0xFE, 0xE9, 0xD6)
RED_LIGHT = RGBColor(0xFB, 0xE2, 0xDF)

FONT_KR = "맑은 고딕"

# ============================================================
# Setup
# ============================================================
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
blank_layout = prs.slide_layouts[6]
TOTAL = 10

# ============================================================
# Helpers
# ============================================================
def add_rect(slide, x, y, w, h, fill=NAVY, line=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    if line is None:
        shape.line.fill.background()
    else:
        shape.line.color.rgb = line
    shape.shadow.inherit = False
    return shape

def add_text(slide, x, y, w, h, text, size=18, bold=False, color=GRAY_TEXT,
             align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, font=FONT_KR):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.05)
    tf.margin_right = Inches(0.05)
    tf.margin_top = Inches(0.02)
    tf.margin_bottom = Inches(0.02)
    tf.vertical_anchor = anchor
    lines = text.split("\n") if isinstance(text, str) else text
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        run = p.add_run()
        run.text = line
        run.font.name = font
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.color.rgb = color
    return tb

def add_notes(slide, narration):
    slide.notes_slide.notes_text_frame.text = narration

def add_page_header(slide, slide_no, section_label, focus_marker=""):
    add_rect(slide, Inches(0), Inches(0), prs.slide_width, Inches(0.12), fill=ORANGE)
    add_text(slide, Inches(0.4), Inches(0.22), Inches(3.0), Inches(0.4),
             "UTTEC", size=14, bold=True, color=NAVY)
    add_text(slide, Inches(0.4), Inches(0.5), Inches(7.0), Inches(0.35),
             "한국기계 귀하 — AI 임직원 교육 7주 커리큘럼 제안", size=10, color=GRAY_MID)
    label = f"{slide_no:02d} / {TOTAL:02d}  ·  {section_label}"
    if focus_marker:
        label = f"{focus_marker}  {label}"
    add_text(slide, Inches(10.5), Inches(0.28), Inches(2.6), Inches(0.4),
             label, size=9, color=GRAY_MID, align=PP_ALIGN.RIGHT)

def add_page_footer(slide):
    add_rect(slide, Inches(0), Inches(7.35), prs.slide_width, Inches(0.15), fill=NAVY)
    add_text(slide, Inches(0.4), Inches(7.37), Inches(8), Inches(0.13),
             "UTTEC · 홍광선  ihong9059@gmail.com · 010-2401-9059 · 경기 용인시 기흥구 흥덕유타워",
             size=8, color=WHITE)
    add_text(slide, Inches(11.5), Inches(7.37), Inches(1.6), Inches(0.13),
             "🔒 회사 고유 기술 외부 노출 0% 보장",
             size=8, color=ORANGE, align=PP_ALIGN.RIGHT)

# ============================================================
# S1 — 타이틀
# ============================================================
s = prs.slides.add_slide(blank_layout)
add_rect(s, Inches(0), Inches(0), prs.slide_width, prs.slide_height, fill=NAVY)
add_rect(s, Inches(0), Inches(0), prs.slide_width, Inches(0.18), fill=ORANGE)
add_text(s, Inches(0.6), Inches(0.4), Inches(5), Inches(0.5),
         "UTTEC × 한국기계", size=18, bold=True, color=ORANGE)
add_text(s, Inches(0.6), Inches(0.95), Inches(7), Inches(0.4),
         "AI 도입 제안 — 임직원 교육 7주 커리큘럼", size=14, color=WHITE)
add_text(s, Inches(0.6), Inches(2.5), Inches(13), Inches(1.4),
         "한국기계 귀하", size=72, bold=True, color=WHITE)
add_text(s, Inches(0.6), Inches(3.9), Inches(13), Inches(0.8),
         "AI 임직원 교육 7주 커리큘럼", size=36, color=ORANGE)
add_text(s, Inches(0.6), Inches(5.0), Inches(13), Inches(0.5),
         "보안 4중 방어  ·  업무효율화 → 국책과제 → 제품 개발", size=22, color=WHITE)

# bottom badge — security focus
add_rect(s, Inches(0.6), Inches(6.2), Inches(12), Inches(0.7),
         fill=GREEN_OK)
add_text(s, Inches(0.6), Inches(6.27), Inches(12), Inches(0.55),
         "🔒  회사 고유 기술 · 외부 노출 0% 보장 — 한국기계 1순위 우려에 대한 답",
         size=18, bold=True, color=WHITE, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

add_notes(s, "안녕하십니까, UTTEC 홍광선입니다. 본 영상은 한국기계 임직원을 위한 AI 도입 교육 7주 커리큘럼 제안입니다. 한국기계가 가장 우려하시는 회사 고유 기술 외부 노출 0%를 보장하는 보안 4중 방어 구조로 설계했습니다.")

# ============================================================
# S2 — 한국기계 우려 → 우리의 답 ⭐
# ============================================================
s = prs.slides.add_slide(blank_layout)
add_page_header(s, 2, "한국기계 1순위 우려", "⭐")
add_page_footer(s)
add_text(s, Inches(0.6), Inches(1.0), Inches(12), Inches(0.6),
         "⭐ 한국기계 1순위 우려에 대한 답", size=30, bold=True, color=NAVY)

# 좌측 — 우려 (Red)
add_rect(s, Inches(0.6), Inches(1.9), Inches(5.7), Inches(4.8),
         fill=RED_LIGHT)
add_rect(s, Inches(0.6), Inches(1.9), Inches(5.7), Inches(0.6), fill=RED_FOCUS)
add_text(s, Inches(0.6), Inches(1.95), Inches(5.7), Inches(0.5),
         "한국기계의 우려", size=20, bold=True, color=WHITE,
         align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
add_text(s, Inches(0.8), Inches(2.7), Inches(5.3), Inches(1.4),
         '"AI 사용하다 보면 우리 회사 고유 기술이\n외부로 노출되지 않을까?"',
         size=18, bold=True, color=RED_FOCUS)
add_text(s, Inches(0.8), Inches(4.3), Inches(5.3), Inches(2.3),
         "•  도면번호  ·  공정 비법\n•  고객 정보  ·  재료 배합\n•  특허·노하우  ·  견적·원가\n•  설계·검사 자료\n\n→ AI에 입력 시 외부 LLM 서버로 전송",
         size=15, color=GRAY_TEXT)

# 화살표
add_text(s, Inches(6.3), Inches(4.0), Inches(0.7), Inches(0.7),
         "→", size=44, bold=True, color=ORANGE, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

# 우측 — 답 (Green)
add_rect(s, Inches(7.05), Inches(1.9), Inches(5.7), Inches(4.8),
         fill=GREEN_LIGHT)
add_rect(s, Inches(7.05), Inches(1.9), Inches(5.7), Inches(0.6), fill=GREEN_OK)
add_text(s, Inches(7.05), Inches(1.95), Inches(5.7), Inches(0.5),
         "UTTEC의 답 — 보안 4중 방어", size=20, bold=True, color=WHITE,
         align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
add_text(s, Inches(7.25), Inches(2.75), Inches(5.4), Inches(3.9),
         "✅  마스킹 SOP (입력 단계)\n     — Phase 1부터\n\n✅  On-Device AI (Ollama 로컬)\n     — 외부 호출 0%  ⭐ Phase 3\n\n✅  사내 저장소 격리\n     — 전 기간 일관\n\n✅  API zero-retention\n     — enterprise 옵션",
         size=15, color=NAVY, bold=True)

add_notes(s, "한국기계가 가장 우려하시는 점을 알고 있습니다. AI를 쓰다 보면 회사의 도면, 고객 정보, 공정 비법, 재료 배합 같은 고유 기술이 외부로 새 나가지 않을까 하는 점이지요. 본 교육은 이 우려에 대해 네 가지 방어층으로 답합니다. 입력 단계의 마스킹 SOP, 외부 호출이 전혀 없는 온디바이스 AI, 사내 저장소 격리, 그리고 API zero-retention 옵션입니다. 첫 강의 시간부터 마지막 강의 시간까지 이 네 가지 방어가 일관되게 작동합니다.")

# ============================================================
# S3 — 3대 지향 + 7주 3 Phase
# ============================================================
s = prs.slides.add_slide(blank_layout)
add_page_header(s, 3, "전체 그림")
add_page_footer(s)
add_text(s, Inches(0.6), Inches(1.0), Inches(12), Inches(0.6),
         "3대 지향 ↔ 7주 3 Phase", size=30, bold=True, color=NAVY)

# 상단 3대 지향
for i, (icon, title, body, c) in enumerate([
    ("🎯", "업무 효율성", "반복 업무 자동화\n문서·보고서 시간 절감", ORANGE),
    ("📚", "기술 확보·정리", "사내 분산 노하우\n→ 단일 지식 베이스", GREEN_OK),
    ("🚀", "비전 구축", "AI 기반 신제품·서비스\n12개월 로드맵", NAVY_LIGHT),
]):
    x = Inches(0.6 + i * 4.15)
    add_rect(s, x, Inches(1.9), Inches(3.85), Inches(1.5), fill=c)
    add_text(s, x, Inches(1.95), Inches(3.85), Inches(0.5),
             f"{icon}  {title}", size=18, bold=True, color=WHITE,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    add_text(s, x, Inches(2.55), Inches(3.85), Inches(0.9), body,
             size=13, color=WHITE, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

# 하단 3 Phase
add_text(s, Inches(0.6), Inches(3.7), Inches(12), Inches(0.5),
         "7주 3 Phase 단계 진행", size=20, bold=True, color=NAVY)
for i, (ph, weeks, title, body, c) in enumerate([
    ("Phase 1", "2주", "AI 기초 + 업무효율화",
     "Claude Code 설치\n마스킹 SOP\nObsidian + Skill", NAVY_LIGHT),
    ("Phase 2", "2주", "국책과제 보고서 자동화",
     "RFP 분석\n사업계획서·중간·결과보고서\n발표자료 일관 워크플로우", GREEN_OK),
    ("Phase 3", "3주", "제품 개발 적용",
     "On-Device AI (Ollama)\n임베디드 AI (TinyML)\n신제품 PoC + 비전", ORANGE),
]):
    x = Inches(0.6 + i * 4.15)
    add_rect(s, x, Inches(4.3), Inches(3.85), Inches(0.7), fill=c)
    add_text(s, x, Inches(4.32), Inches(3.85), Inches(0.65),
             f"{ph}  ·  {weeks}", size=17, bold=True, color=WHITE,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    add_rect(s, x, Inches(5.0), Inches(3.85), Inches(0.55), fill=GRAY_LIGHT)
    add_text(s, x, Inches(5.05), Inches(3.85), Inches(0.5),
             title, size=14, bold=True, color=NAVY,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    add_rect(s, x, Inches(5.55), Inches(3.85), Inches(1.5), fill=CREAM)
    add_text(s, x.emu + Inches(0.15).emu, Inches(5.6).emu, Inches(3.55), Inches(1.4), body,
             size=12, color=GRAY_TEXT)

add_text(s, Inches(0.6), Inches(7.1), Inches(12), Inches(0.25),
         "총 7주 = 14 강의 일 (한 주 2일 4시간 × 7)",
         size=12, color=GRAY_MID, align=PP_ALIGN.CENTER)

add_notes(s, "교육 목적은 사용자께서 명시하신 세 가지입니다. 첫째, 임직원 업무 효율성. 둘째, 회사 보유 기술의 확보와 단일 지식 베이스화. 셋째, AI 기반 향후 회사 비전 구축입니다. 이 세 가지 지향을 7주 3 Phase로 단계화했습니다. Phase 1 2주, AI 기초와 업무 효율화. Phase 2 2주, 국책과제 보고서 실무 적용. Phase 3 3주, 온디바이스 AI로 제품 개발 적용입니다. 한 주 2일 4시간, 총 14 강의 일입니다.")

# ============================================================
# S4 — Phase 1
# ============================================================
s = prs.slides.add_slide(blank_layout)
add_page_header(s, 4, "Phase 1 — 기초")
add_page_footer(s)
add_rect(s, Inches(0.6), Inches(0.95), Inches(12), Inches(0.55), fill=NAVY_LIGHT)
add_text(s, Inches(0.6), Inches(0.95), Inches(12), Inches(0.55),
         "Phase 1 — AI 기초 + 업무 효율화 (2주)  ⭐ 보안 첫날부터",
         size=20, bold=True, color=WHITE, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

# W1
add_rect(s, Inches(0.6), Inches(1.8), Inches(5.95), Inches(4.4), fill=CREAM)
add_rect(s, Inches(0.6), Inches(1.8), Inches(5.95), Inches(0.55), fill=ORANGE)
add_text(s, Inches(0.6), Inches(1.83), Inches(5.95), Inches(0.5),
         "W1 — AI 첫 만남 + 마스킹 SOP", size=16, bold=True, color=WHITE,
         align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
add_text(s, Inches(0.8), Inches(2.5), Inches(5.55), Inches(3.6),
         "Day 1 (4h)\n  •  AI 시대 개요 + 한국 도입 사례\n  •  ⭐ 마스킹 SOP — 도면·고객·공정\n  •  마스킹 실습 ([CUST-A] [PROC-001])\n  •  Claude Code 설치\n\nDay 2 (4h)\n  •  Claude Code 첫 대화\n  •  위험 시연 (마스킹 안 된 입력)\n  •  안전 시연 (마스킹 입력 = 동일 결과)\n  •  사내 마스킹 정책 초안 작성",
         size=13, color=GRAY_TEXT)

# W2
add_rect(s, Inches(6.78), Inches(1.8), Inches(5.95), Inches(4.4), fill=CREAM)
add_rect(s, Inches(6.78), Inches(1.8), Inches(5.95), Inches(0.55), fill=GREEN_OK)
add_text(s, Inches(6.78), Inches(1.83), Inches(5.95), Inches(0.5),
         "W2 — Obsidian + Skill 자동화", size=16, bold=True, color=WHITE,
         align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
add_text(s, Inches(6.98), Inches(2.5), Inches(5.55), Inches(3.6),
         "Day 1 (4h)\n  •  Obsidian 설치 + 한국기계_지식베이스 Vault\n  •  마크다운 + 양방향 링크 [[페이지명]]\n  •  사내 카테고리 (장비·공정·고객·gotcha)\n  •  Claude + Obsidian 연계\n\nDay 2 (4h)\n  •  Skill 개념 (반복 업무 → 슬래시 명령)\n  •  /work-start / /work-end 설치\n  •  /일일보고서 (사내 결재 양식)\n  •  MCP — 사내 Notion 연결",
         size=13, color=GRAY_TEXT)

# 산출물
add_rect(s, Inches(0.6), Inches(6.35), Inches(12.15), Inches(0.7), fill=GREEN_LIGHT)
add_text(s, Inches(0.6), Inches(6.4), Inches(12.15), Inches(0.6),
         "산출물:  📄 마스킹 정책 v1.0    📚 지식베이스 100 페이지    ⚙️ 사내 Skill 3종",
         size=15, bold=True, color=NAVY, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

add_notes(s, "Phase 1은 AI를 도구로 인식하게 만드는 단계입니다. 첫날부터 마스킹 SOP를 가르칩니다. 회사 문서를 AI에 입력하기 전에 도면번호, 고객명, 공정 비법을 placeholder로 치환하는 표준 절차입니다. 실제 위험 시연으로 마스킹 안 된 문서가 외부로 어떻게 전송되는지 보여드리고, 마스킹된 문서로 같은 결과를 얻을 수 있다는 점도 비교 시연합니다. 2주차는 Obsidian으로 사내에 분산된 노하우를 단일 지식 베이스로 묶고, 반복 업무를 슬래시 명령으로 자동화하는 Skill을 익힙니다.")

# ============================================================
# S5 — Phase 2
# ============================================================
s = prs.slides.add_slide(blank_layout)
add_page_header(s, 5, "Phase 2 — 보고서")
add_page_footer(s)
add_rect(s, Inches(0.6), Inches(0.95), Inches(12), Inches(0.55), fill=GREEN_OK)
add_text(s, Inches(0.6), Inches(0.95), Inches(12), Inches(0.55),
         "Phase 2 — 국책과제 보고서 자동화 (2주)",
         size=20, bold=True, color=WHITE, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

# 좌측 — 4 step 흐름
add_rect(s, Inches(0.6), Inches(1.8), Inches(5.95), Inches(4.4), fill=CREAM)
add_text(s, Inches(0.8), Inches(2.0), Inches(5.55), Inches(0.5),
         "4 step 흐름", size=18, bold=True, color=NAVY)
for i, (step, body) in enumerate([
    ("① RFP 흡수", "RFP 자동 분석 — 핵심 요구사항·평가·예산·기간"),
    ("② 매칭 분석", "자사 자산 ↔ RFP 요구사항 매트릭스 자동"),
    ("③ 사업계획서", "15장 standard /사업계획서_초안 Skill"),
    ("④ 중간/결과", "실험노트·미팅 메모 → 보고서 초안 + 마스킹"),
]):
    y = Inches(2.65 + i * 0.85)
    add_rect(s, Inches(0.8), y, Inches(0.4), Inches(0.7), fill=GREEN_OK)
    add_text(s, Inches(0.8), y, Inches(0.4), Inches(0.7),
             str(i+1), size=18, bold=True, color=WHITE,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    add_text(s, Inches(1.3), y, Inches(5.0), Inches(0.35),
             step, size=15, bold=True, color=NAVY)
    add_text(s, Inches(1.3), y.emu + Inches(0.35).emu, Inches(5.0), Inches(0.35),
             body, size=11, color=GRAY_TEXT)

# 우측 — 정량 효과 매트릭스
add_rect(s, Inches(6.78), Inches(1.8), Inches(5.95), Inches(4.4), fill=GREEN_LIGHT)
add_text(s, Inches(6.98), Inches(2.0), Inches(5.55), Inches(0.5),
         "정량 효과 매트릭스", size=18, bold=True, color=NAVY)

# 테이블 headers
hy = 2.65
add_rect(s, Inches(6.98), Inches(hy), Inches(2.6), Inches(0.45), fill=NAVY)
add_rect(s, Inches(9.58), Inches(hy), Inches(1.4), Inches(0.45), fill=NAVY)
add_rect(s, Inches(10.98), Inches(hy), Inches(1.55), Inches(0.45), fill=NAVY)
add_text(s, Inches(6.98), Inches(hy), Inches(2.6), Inches(0.45),
         "항목", size=13, bold=True, color=WHITE, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
add_text(s, Inches(9.58), Inches(hy), Inches(1.4), Inches(0.45),
         "기존", size=13, bold=True, color=WHITE, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
add_text(s, Inches(10.98), Inches(hy), Inches(1.55), Inches(0.45),
         "AI 적용", size=13, bold=True, color=WHITE, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

rows = [
    ("RFP 분석", "1~2일", "30분"),
    ("사업계획서", "1주", "1~2일"),
    ("중간보고서", "3일", "4시간"),
    ("결과보고서", "1주", "2일"),
    ("발표자료", "3일", "4시간"),
]
for i, (a, b, c) in enumerate(rows):
    ry = hy + 0.45 + i * 0.4
    bg = WHITE if i % 2 == 0 else GRAY_LIGHT
    add_rect(s, Inches(6.98), Inches(ry), Inches(2.6), Inches(0.4), fill=bg)
    add_rect(s, Inches(9.58), Inches(ry), Inches(1.4), Inches(0.4), fill=bg)
    add_rect(s, Inches(10.98), Inches(ry), Inches(1.55), Inches(0.4), fill=bg)
    add_text(s, Inches(7.1), Inches(ry), Inches(2.5), Inches(0.4),
             a, size=12, color=GRAY_TEXT, anchor=MSO_ANCHOR.MIDDLE)
    add_text(s, Inches(9.58), Inches(ry), Inches(1.4), Inches(0.4),
             b, size=12, color=GRAY_TEXT, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    add_text(s, Inches(10.98), Inches(ry), Inches(1.55), Inches(0.4),
             c, size=12, bold=True, color=GREEN_OK, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

# 하단 강조
add_rect(s, Inches(0.6), Inches(6.35), Inches(12.15), Inches(0.7), fill=ORANGE)
add_text(s, Inches(0.6), Inches(6.4), Inches(12.15), Inches(0.6),
         "⭐  국책과제 1건당 약 2주 시간 절감 → 연간 수주 가능 건수 ↑",
         size=18, bold=True, color=WHITE, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

add_notes(s, "Phase 2는 한국기계가 자주 진행하는 국책과제 보고서를 AI로 자동화합니다. RFP 분석부터 사업계획서, 중간보고서, 결과보고서, 그리고 발표 자료까지 한 흐름으로 처리합니다. 정량 효과를 보여드립니다. RFP 분석은 1~2일이 30분으로, 사업계획서 초안은 1주가 1~2일로, 중간보고서는 3일이 4시간으로 단축됩니다. 국책과제 한 건당 약 2주의 시간이 절감됩니다. 이 과정 내내 회사 고유 기술은 마스킹과 검토 게이트로 보호됩니다.")

# ============================================================
# S6 — Phase 3: On-Device AI ⭐⭐⭐
# ============================================================
s = prs.slides.add_slide(blank_layout)
add_page_header(s, 6, "Phase 3 — On-Device", "⭐⭐⭐")
add_page_footer(s)
add_rect(s, Inches(0.6), Inches(0.95), Inches(12), Inches(0.55), fill=ORANGE)
add_text(s, Inches(0.6), Inches(0.95), Inches(12), Inches(0.55),
         "⭐⭐⭐ Phase 3 — On-Device AI (외부 노출 0%)  ·  한국기계 결정타",
         size=20, bold=True, color=WHITE, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

# 좌측 Before (gray)
add_rect(s, Inches(0.6), Inches(1.8), Inches(5.7), Inches(3.5), fill=GRAY_LIGHT)
add_text(s, Inches(0.6), Inches(1.9), Inches(5.7), Inches(0.5),
         "Before  ·  외부 LLM 호출", size=18, bold=True, color=GRAY_MID,
         align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
add_text(s, Inches(0.8), Inches(2.7), Inches(5.3), Inches(2.5),
         "💻 회사 노트북\n      ↓ (인터넷)\n☁️ Claude / GPT 외부 서버\n      ↓\n📤 회사 데이터 전송\n\n→ 마스킹으로 차단 (Phase 1)",
         size=15, color=GRAY_TEXT, align=PP_ALIGN.CENTER)

# 화살표
add_text(s, Inches(6.3), Inches(3.3), Inches(0.7), Inches(0.7),
         "→", size=44, bold=True, color=ORANGE, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

# 우측 After (green)
add_rect(s, Inches(7.05), Inches(1.8), Inches(5.7), Inches(3.5), fill=GREEN_LIGHT)
add_rect(s, Inches(7.05), Inches(1.8), Inches(5.7), Inches(0.55), fill=GREEN_OK)
add_text(s, Inches(7.05), Inches(1.83), Inches(5.7), Inches(0.5),
         "After  ·  Ollama 로컬 SLM", size=18, bold=True, color=WHITE,
         align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
add_text(s, Inches(7.25), Inches(2.55), Inches(5.3), Inches(2.7),
         "💻 회사 노트북 (단독)\n      ↓\n🤖 Ollama / Mistral / Llama\n      ↓\n📁 노트북 안에서 완결\n\n✅ 외부 망 차단 환경에서도 작동",
         size=15, bold=True, color=NAVY, align=PP_ALIGN.CENTER)

# 중앙 큰 강조
add_rect(s, Inches(0.6), Inches(5.45), Inches(12.15), Inches(0.9), fill=NAVY)
add_text(s, Inches(0.6), Inches(5.5), Inches(12.15), Inches(0.8),
         "🔒  외부 망 차단 환경에서도 작동  ·  비공개 데이터 외부 전송 0%",
         size=22, bold=True, color=ORANGE, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

# 하단 W5 디테일
add_text(s, Inches(0.6), Inches(6.5), Inches(12.15), Inches(0.5),
         "W5: Ollama 설치 + Llama 3.1 8B / Mistral 7B  ·  노트북 사양 검증 (RAM 16GB)  ·  망 끊고 작동 검증",
         size=12, color=GRAY_TEXT, align=PP_ALIGN.CENTER)

add_notes(s, "Phase 3가 본 교육의 결정타입니다. 온디바이스 AI라는 기술로, 외부 LLM 호출이 전혀 없이 회사 노트북 안에서 모든 AI 처리를 완결합니다. 5주차에 Ollama 또는 LM Studio를 설치하고, Llama 또는 Mistral 같은 로컬 모델을 노트북에 다운로드합니다. 그리고 외부 망을 완전히 차단한 상태에서 작동하는지 직접 시연으로 확인합니다. 회사 비공개 도면, 고객 데이터, 공정 비법 모두 사내 노트북 밖으로 한 비트도 나가지 않습니다. 이것이 한국기계 우려에 대한 가장 강한 답입니다.")

# ============================================================
# S7 — 임베디드 AI + 제품 통합
# ============================================================
s = prs.slides.add_slide(blank_layout)
add_page_header(s, 7, "Phase 3 — 제품 통합")
add_page_footer(s)
add_rect(s, Inches(0.6), Inches(0.95), Inches(12), Inches(0.55), fill=NAVY_LIGHT)
add_text(s, Inches(0.6), Inches(0.95), Inches(12), Inches(0.55),
         "Phase 3 후반 — 임베디드 AI + 제품 통합 (W6~W7)",
         size=20, bold=True, color=WHITE, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

# 좌측 보드 매트릭스
add_rect(s, Inches(0.6), Inches(1.8), Inches(3.9), Inches(4.4), fill=CREAM)
add_text(s, Inches(0.8), Inches(2.0), Inches(3.5), Inches(0.5),
         "UTTEC 검증 자산", size=16, bold=True, color=NAVY)
add_text(s, Inches(0.8), Inches(2.6), Inches(3.5), Inches(3.4),
         "📋  14 보드 매트릭스\n     Cortex-M4F / M7 / LX7 / RISC-V\n\n🔬  Round 1~50 검증 carry\n     R18 3.23× / R44 KWS 9.91ms\n     R46 9.26ms / R50 8.13ms\n\n📦  sensor 라이브러리 12 모듈\n     MPU / BME680 / FT5336\n     INMP441 / INA219 / 외 7\n\n📊  AI 매트릭스 491 lines\n     11 sensor × 2~3 AI model\n     × 14 보드 최저선",
         size=11, color=GRAY_TEXT)

# 중앙 워크플로우
add_rect(s, Inches(4.65), Inches(1.8), Inches(3.5), Inches(4.4), fill=GREEN_LIGHT)
add_text(s, Inches(4.85), Inches(2.0), Inches(3.1), Inches(0.5),
         "워크플로우 §0.4 (6 step)", size=16, bold=True, color=NAVY,
         align=PP_ALIGN.CENTER)
for i, st in enumerate([
    "① 사내 sensor 데이터 수집\n      (외부 데이터 의존 0%)",
    "② PC 학습 (PyTorch)",
    "③ PC 검증 (sanity)",
    "④ INT8 quantize",
    "⑤ MCU port (CMSIS-NN)",
    "⑥ MCU sweep (3축 검증)",
]):
    y = Inches(2.65 + i * 0.55)
    add_text(s, Inches(4.85), y, Inches(3.1), Inches(0.5),
             st, size=10, color=GRAY_TEXT, anchor=MSO_ANCHOR.MIDDLE)

# 우측 한국기계 응용
add_rect(s, Inches(8.3), Inches(1.8), Inches(4.45), Inches(4.4), fill=ORANGE_LIGHT)
add_text(s, Inches(8.5), Inches(2.0), Inches(4.05), Inches(0.5),
         "한국기계 응용 brainstorm", size=16, bold=True, color=NAVY)
add_text(s, Inches(8.5), Inches(2.6), Inches(4.05), Inches(3.4),
         "🔧  예측 정비\n     (vibration · current 이상 감지)\n\n🔍  품질 검사\n     (vision · 치수 측정)\n\n⚠️  이상 감지\n     (온도·소리·전류 패턴)\n\n📈  공정 최적화\n     (실시간 sensor → 자동 조정)\n\n→  6주차 PoC 1건 선정\n→  7주차 데이터 → 모델 → 시연",
         size=12, color=GRAY_TEXT)

# 하단 강조
add_rect(s, Inches(0.6), Inches(6.35), Inches(12.15), Inches(0.7), fill=ORANGE)
add_text(s, Inches(0.6), Inches(6.4), Inches(12.15), Inches(0.6),
         "신제품 PoC 1건 → 양산 진입 시 차별화 + AI 비전 로드맵 12개월",
         size=17, bold=True, color=WHITE, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

add_notes(s, "Phase 3 후반은 한국기계 제품에 AI를 직접 넣는 단계입니다. 임베디드 AI 또는 TinyML이라는 기술로, sensor가 달린 제품 안에서 직접 AI 추론이 돌아갑니다. UTTEC가 50회 라운드로 검증한 14개 보드 매트릭스와 12개 sensor 라이브러리를 그대로 활용합니다. 한국기계 제품에 어떤 응용이 가능한지 brainstorm으로 시작합니다. 예측 정비, 품질 검사, 이상 감지 같은 영역입니다. 6주차에 PoC 1건을 선정하고 7주차에 실제 데이터로 시연합니다.")

# ============================================================
# S8 — 보안 4중 방어 도식 ⭐⭐⭐
# ============================================================
s = prs.slides.add_slide(blank_layout)
add_page_header(s, 8, "보안 핵심", "⭐⭐⭐")
add_page_footer(s)
add_rect(s, Inches(0.6), Inches(0.95), Inches(12), Inches(0.7), fill=ORANGE)
add_text(s, Inches(0.6), Inches(0.95), Inches(12), Inches(0.7),
         "⭐⭐⭐  보안 4중 방어  ·  회사 고유 기술 외부 노출 0%",
         size=22, bold=True, color=WHITE, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

# 4 컬럼
cols = [
    ("1", "마스킹 SOP", "입력 단계", "Phase 1부터", NAVY_LIGHT),
    ("2", "On-Device AI", "외부 호출 0%", "⭐ Phase 3 결정타", ORANGE),
    ("3", "사내 저장소 격리", "외부 동기화 금지", "전 기간 일관", GREEN_OK),
    ("4", "API zero-retention", "enterprise 옵션", "Phase 2~3 부득이 시", NAVY_LIGHT),
]
for i, (n, t, sub, ph, c) in enumerate(cols):
    x = Inches(0.6 + i * 3.1)
    add_rect(s, x, Inches(2.0), Inches(2.85), Inches(4.0), fill=CREAM)
    # number circle
    add_rect(s, x.emu + Inches(0.9).emu, Inches(2.2).emu, Inches(1.05), Inches(1.05),
             fill=c)
    add_text(s, x.emu + Inches(0.9).emu, Inches(2.2).emu, Inches(1.05), Inches(1.05),
             n, size=44, bold=True, color=WHITE,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    add_text(s, x.emu + Inches(0.15).emu, Inches(3.5).emu, Inches(2.55), Inches(0.5),
             f"🔒  {t}", size=16, bold=True, color=NAVY,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    add_text(s, x.emu + Inches(0.15).emu, Inches(4.1).emu, Inches(2.55), Inches(0.4),
             sub, size=13, color=GRAY_TEXT, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    add_text(s, x.emu + Inches(0.15).emu, Inches(4.6).emu, Inches(2.55), Inches(0.4),
             ph, size=12, bold=True, color=ORANGE, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

# 하단 강조
add_rect(s, Inches(0.6), Inches(6.35), Inches(12.15), Inches(0.7), fill=NAVY)
add_text(s, Inches(0.6), Inches(6.4), Inches(12.15), Inches(0.6),
         "한국기계가 가장 우려하신 점에 대한 우리의 답",
         size=20, bold=True, color=ORANGE, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

add_notes(s, "보안 4중 방어를 다시 한번 정리해 드립니다. 첫째, 입력 단계 마스킹 SOP — Phase 1부터. 둘째, 온디바이스 AI로 외부 호출 0% — Phase 3 결정타. 셋째, 사내 저장소 격리 — 전 기간. 넷째, 부득이한 외부 API 사용 시 zero-retention enterprise 옵션. 네 가지 방어가 일관되게 작동합니다. 한국기계 1순위 우려에 대한 우리의 답입니다.")

# ============================================================
# S9 — 효과 + 비용
# ============================================================
s = prs.slides.add_slide(blank_layout)
add_page_header(s, 9, "효과 + 비용")
add_page_footer(s)
add_rect(s, Inches(0.6), Inches(0.95), Inches(12), Inches(0.55), fill=GREEN_OK)
add_text(s, Inches(0.6), Inches(0.95), Inches(12), Inches(0.55),
         "효과 + 비용",
         size=22, bold=True, color=WHITE, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

# 좌측 효과
add_rect(s, Inches(0.6), Inches(1.8), Inches(7.5), Inches(4.4), fill=GREEN_LIGHT)
add_text(s, Inches(0.8), Inches(2.0), Inches(7.1), Inches(0.5),
         "정량 + 정성 효과", size=18, bold=True, color=NAVY)
add_text(s, Inches(0.8), Inches(2.7), Inches(7.1), Inches(3.5),
         "✅  국책과제 1건당 약 2주 시간 절감\n      → 연간 수주 가능 건수 ↑\n\n✅  일상 보고서·문서 50% 시간 단축\n\n✅  회사 고유 기술 외부 노출 0%  ⭐ 결정타\n\n✅  사내 분산 노하우 → 단일 지식 베이스\n      (Obsidian Vault 영구 자산)\n\n✅  신제품 PoC 1건  → 양산 차별화\n\n✅  임직원 AI 친화도 ↑  → 기술 영입 attractive",
         size=14, color=GRAY_TEXT)

# 우측 비용
add_rect(s, Inches(8.3), Inches(1.8), Inches(4.45), Inches(4.4), fill=CREAM)
add_text(s, Inches(8.5), Inches(2.0), Inches(4.05), Inches(0.5),
         "비용 견적", size=18, bold=True, color=NAVY)

# 테이블
hy = 2.65
add_rect(s, Inches(8.5), Inches(hy), Inches(2.6), Inches(0.4), fill=NAVY)
add_rect(s, Inches(11.1), Inches(hy), Inches(1.45), Inches(0.4), fill=NAVY)
add_text(s, Inches(8.5), Inches(hy), Inches(2.6), Inches(0.4),
         "항목", size=12, bold=True, color=WHITE, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
add_text(s, Inches(11.1), Inches(hy), Inches(1.45), Inches(0.4),
         "견적", size=12, bold=True, color=WHITE, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
brows = [
    ("강사료 (14일 × 100만)", "1,400만"),
    ("TinyML 보드·sensor", "200~300만"),
]
for i, (a, b) in enumerate(brows):
    ry = hy + 0.4 + i * 0.45
    bg = WHITE if i % 2 == 0 else GRAY_LIGHT
    add_rect(s, Inches(8.5), Inches(ry), Inches(2.6), Inches(0.45), fill=bg)
    add_rect(s, Inches(11.1), Inches(ry), Inches(1.45), Inches(0.45), fill=bg)
    add_text(s, Inches(8.6), Inches(ry), Inches(2.5), Inches(0.45),
             a, size=11, color=GRAY_TEXT, anchor=MSO_ANCHOR.MIDDLE)
    add_text(s, Inches(11.1), Inches(ry), Inches(1.45), Inches(0.45),
             b, size=11, color=GRAY_TEXT, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
# 합계
ry = hy + 0.4 + 2 * 0.45
add_rect(s, Inches(8.5), Inches(ry), Inches(2.6), Inches(0.55), fill=ORANGE)
add_rect(s, Inches(11.1), Inches(ry), Inches(1.45), Inches(0.55), fill=ORANGE)
add_text(s, Inches(8.6), Inches(ry), Inches(2.5), Inches(0.55),
         "합계", size=14, bold=True, color=WHITE, anchor=MSO_ANCHOR.MIDDLE)
add_text(s, Inches(11.1), Inches(ry), Inches(1.45), Inches(0.55),
         "1,600~1,700만", size=12, bold=True, color=WHITE, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

# 정부지원 안내
add_rect(s, Inches(8.5), Inches(5.0), Inches(4.05), Inches(1.15), fill=GREEN_LIGHT)
add_text(s, Inches(8.6), Inches(5.1), Inches(3.95), Inches(0.4),
         "💡 정부지원 활용 시", size=13, bold=True, color=NAVY)
add_text(s, Inches(8.6), Inches(5.5), Inches(3.95), Inches(0.6),
         "중기부 디지털 인재 양성\n70~80% 환급 가능 (신청 후 6개월 lead)",
         size=11, color=GRAY_TEXT)

# 하단
add_text(s, Inches(0.6), Inches(6.4), Inches(12.15), Inches(0.5),
         "ROI: 국책과제 1건 절감 + 신제품 PoC 차별화 = 1~2년 내 회수",
         size=14, bold=True, color=NAVY, align=PP_ALIGN.CENTER)

add_notes(s, "효과와 비용을 정리해드립니다. 정량 효과는 국책과제 1건당 약 2주 시간 절감, 일상 보고서 50% 단축입니다. 정성 효과는 회사 고유 기술 외부 노출 0%, 사내 노하우 단일 지식 베이스화, 그리고 신제품 PoC 한 건의 양산 차별화입니다. 비용은 강사료 1,400만 원, TinyML 보드 200~300만 원, 합계 약 1,600~1,700만 원입니다. 중기부 디지털 인재 양성 사업을 활용하시면 70~80% 환급이 가능합니다. 신청 후 6개월 lead가 필요합니다.")

# ============================================================
# S10 — 다음 단계 + 클로징
# ============================================================
s = prs.slides.add_slide(blank_layout)
add_rect(s, Inches(0), Inches(0), prs.slide_width, prs.slide_height, fill=NAVY)
add_rect(s, Inches(0), Inches(0), prs.slide_width, Inches(0.18), fill=ORANGE)
add_text(s, Inches(0.6), Inches(0.5), Inches(12), Inches(0.5),
         "다음 단계", size=28, bold=True, color=ORANGE)

# 3 step
steps = [
    ("1", "본 영상 + 커리큘럼 검토", "한국기계 측 — 내용·일정·예산"),
    ("2", "사전 미팅", "수강생 명단 · 일정 · 노트북 사양"),
    ("3", "교육 시작", "2주 전 사전 환경 점검 → 첫 강의"),
]
for i, (n, t, body) in enumerate(steps):
    x = Inches(0.6 + i * 4.15)
    add_rect(s, x, Inches(1.5), Inches(3.85), Inches(2.0), fill=NAVY_LIGHT)
    add_rect(s, x.emu + Inches(1.4).emu, Inches(1.7).emu, Inches(1.05), Inches(1.05), fill=ORANGE)
    add_text(s, x.emu + Inches(1.4).emu, Inches(1.7).emu, Inches(1.05), Inches(1.05),
             n, size=44, bold=True, color=WHITE,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    add_text(s, x, Inches(2.9), Inches(3.85), Inches(0.35),
             t, size=15, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text(s, x, Inches(3.25), Inches(3.85), Inches(0.3),
             body, size=11, color=GRAY_LIGHT, align=PP_ALIGN.CENTER)

# 중앙 큰 메시지
add_text(s, Inches(0.6), Inches(4.0), Inches(12.15), Inches(1.0),
         "한국기계의 AI 도입",
         size=44, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
add_text(s, Inches(0.6), Inches(4.95), Inches(12.15), Inches(0.6),
         "보안  ·  효율  ·  비전  ·  함께 만들어 가겠습니다",
         size=22, color=ORANGE, align=PP_ALIGN.CENTER)

# 보안 강조
add_rect(s, Inches(0.6), Inches(5.85), Inches(12.15), Inches(0.7), fill=GREEN_OK)
add_text(s, Inches(0.6), Inches(5.9), Inches(12.15), Inches(0.6),
         "🔒  회사 고유 기술 · 외부 노출 0% 보장",
         size=20, bold=True, color=WHITE, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

# 푸터
add_text(s, Inches(0.6), Inches(6.8), Inches(12.15), Inches(0.3),
         "UTTEC · 홍광선  ihong9059@gmail.com · 010-2401-9059",
         size=14, color=WHITE, align=PP_ALIGN.CENTER)
add_text(s, Inches(0.6), Inches(7.15), Inches(12.15), Inches(0.25),
         "본 영상 + 커리큘럼.md + 슬라이드.pptx 전부 제공",
         size=11, color=GRAY_LIGHT, align=PP_ALIGN.CENTER)

add_notes(s, "다음 단계는 세 가지입니다. 첫째, 본 영상과 커리큘럼을 한국기계 측에서 검토하십시오. 둘째, 사전 미팅에서 수강생 명단, 일정, 노트북 사양을 정합니다. 셋째, 교육 시작 2주 전 사전 환경 점검을 진행합니다. 한국기계의 AI 도입을 보안, 효율, 비전 세 가지 축으로 함께 만들어 가겠습니다. 감사합니다.")

# ============================================================
# Save
# ============================================================
import os
out = os.path.join(os.path.dirname(os.path.abspath(__file__)),
                   "한국기계_AI교육_제안.pptx")
prs.save(out)
print(f"saved: {out}")
print(f"slides: {len(prs.slides)} / target 10")

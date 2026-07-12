# -*- coding: utf-8 -*-
"""
UTTEC 골프장 LoRa 무선 통합 제어 시스템 제안서 (범용본) 생성
- 2대 시스템: 수조 관리 + 야간 조명
- 견적 = 규모별 프레임만(구체 단가 없음)
- 야간조명 = 설계 적용 톤 (무선조명 8년 상용화 실적 + 순차소등은 신규 설계)
- 도식 기반 (지형도면은 별도 삽입 가능)
- v2(2026-07-11): 슬라이드 6 '수위 변화 추적·누수 조기 진단'(부표식 대비) 신설 → 18슬라이드
출력: UTTEC_LoRa무선제어_제안서_범용_v2_20260711.pptx
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR

# ---- 색상 팔레트 ----
NAVY   = RGBColor(0x0F, 0x23, 0x40)   # 제목/본문 진한 남색
ORANGE = RGBColor(0xEA, 0x58, 0x0C)   # UTTEC 액센트
GRAY   = RGBColor(0x4B, 0x55, 0x63)   # 부제/본문
LIGHT  = RGBColor(0xF1, 0xF5, 0xF9)   # 밝은 배경 박스
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
GREEN  = RGBColor(0x15, 0x80, 0x3D)
RED    = RGBColor(0xDC, 0x2E, 0x26)
FONT   = "맑은 고딕"

prs = Presentation()
prs.slide_width  = Inches(13.333)   # 16:9
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]
SW, SH = prs.slide_width, prs.slide_height

def add_slide():
    return prs.slides.add_slide(BLANK)

def _set(tf, size, color, bold=False, align=PP_ALIGN.LEFT):
    for p in tf.paragraphs:
        p.alignment = align
        for r in p.runs:
            r.font.name = FONT; r.font.size = Pt(size)
            r.font.color.rgb = color; r.font.bold = bold

def box(slide, l, t, w, h, fill=None, line=None):
    sp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, l, t, w, h)
    sp.shadow.inherit = False
    if fill is None: sp.fill.background()
    else: sp.fill.solid(); sp.fill.fore_color.rgb = fill
    if line is None: sp.line.fill.background()
    else: sp.line.color.rgb = line; sp.line.width = Pt(1)
    return sp

def text(slide, l, t, w, h, s, size=18, color=NAVY, bold=False,
         align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP):
    tb = slide.shapes.add_textbox(l, t, w, h); tf = tb.text_frame
    tf.word_wrap = True; tf.vertical_anchor = anchor
    lines = s.split("\n")
    tf.text = lines[0]
    for extra in lines[1:]:
        p = tf.add_paragraph(); p.text = extra
    _set(tf, size, color, bold, align)
    return tb

def bullets(slide, l, t, w, h, items, size=17, gap=6):
    tb = slide.shapes.add_textbox(l, t, w, h); tf = tb.text_frame
    tf.word_wrap = True
    for i, (txt, lv, col, bd) in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.level = lv; p.space_after = Pt(gap)
        run = p.add_run(); run.text = ("• " if lv == 0 else "– ") + txt
        run.font.name = FONT; run.font.size = Pt(size - lv*1)
        run.font.color.rgb = col; run.font.bold = bd
    return tb

def header(slide, no, title, sub=None):
    box(slide, 0, 0, SW, Inches(1.15), fill=NAVY)
    box(slide, 0, Inches(1.15), SW, Pt(4), fill=ORANGE)
    text(slide, Inches(0.5), Inches(0.12), Inches(11.8), Inches(0.55),
         title, size=26, color=WHITE, bold=True, anchor=MSO_ANCHOR.MIDDLE)
    if sub:
        text(slide, Inches(0.52), Inches(0.72), Inches(11.8), Inches(0.35),
             sub, size=13, color=RGBColor(0xCB,0xD5,0xE1))
    text(slide, Inches(12.4), Inches(0.35), Inches(0.7), Inches(0.5),
         f"{no:02d}", size=20, color=ORANGE, bold=True, align=PP_ALIGN.RIGHT)

def footer(slide):
    text(slide, Inches(0.5), Inches(7.05), Inches(8), Inches(0.35),
         "UTTEC  |  LoRa 무선 통합 제어 시스템 제안", size=10, color=GRAY)

def stat(slide, l, t, w, label, value, sub, vcolor=ORANGE):
    box(slide, l, t, w, Inches(1.75), fill=LIGHT)
    box(slide, l, t, w, Pt(5), fill=vcolor)
    text(slide, l, t+Inches(0.18), w, Inches(0.45), label, size=15, color=GRAY, align=PP_ALIGN.CENTER)
    text(slide, l, t+Inches(0.6), w, Inches(0.7), value, size=32, color=vcolor, bold=True, align=PP_ALIGN.CENTER)
    text(slide, l, t+Inches(1.28), w, Inches(0.4), sub, size=12, color=GRAY, align=PP_ALIGN.CENTER)

# ================= Slide 1 — 표지 =================
s = add_slide()
box(s, 0, 0, SW, SH, fill=NAVY)
box(s, 0, Inches(4.15), SW, Pt(5), fill=ORANGE)
text(s, Inches(0.9), Inches(1.5), Inches(11.5), Inches(0.6),
     "골프장 시설관리부서 귀중", size=18, color=RGBColor(0xCB,0xD5,0xE1))
text(s, Inches(0.9), Inches(2.3), Inches(11.5), Inches(1.6),
     "골프장 LoRa 무선 통합 제어 시스템", size=44, color=WHITE, bold=True)
text(s, Inches(0.9), Inches(3.5), Inches(11.5), Inches(0.7),
     "수조 자동 관리  +  야간 조명 순차 제어", size=26, color=ORANGE, bold=True)
text(s, Inches(0.9), Inches(4.5), Inches(11.5), Inches(0.6),
     "전력비 절감 · 무인 운영 · 인적 사고 예방", size=18, color=RGBColor(0xE2,0xE8,0xF0))
text(s, Inches(0.9), Inches(6.3), Inches(11.5), Inches(0.8),
     "UTTEC  ·  국내 최초 LoRa 무선 조명 상용화(2018)  ·  제안일 2026-07-05",
     size=14, color=RGBColor(0x94,0xA3,0xB8))

# ================= Slide 2 — 운영 현장의 고민 =================
s = add_slide(); header(s, 2, "운영 현장의 고민", "골프장 시설관리에서 반복되는 4가지 부담")
cards = [
    ("야간 조명 수동 소등", "마지막 팀이 지나간 홀 조명을 끄러\n담당자가 밤에 직접 이동"),
    ("수조 수위 수동 점검", "고지대 물탱크 수위를 매번 확인하고\n원거리 펌프를 수동 조작"),
    ("위험 위치 야간 조작", "경사지·외곽의 차단기·스위치를\n야간·우천에 직접 조작 (사고 위험)"),
    ("노후 설비 긴급 AS", "설치 5년+ 부품 노화기\n고장 시 즉시 영업 차질"),
]
x0 = Inches(0.5); y0 = Inches(1.6); cw = Inches(6.0); ch = Inches(2.4); gx = Inches(0.33); gy = Inches(0.35)
for i,(t_,b_) in enumerate(cards):
    cx = x0 + (cw+gx)*(i%2); cy = y0 + (ch+gy)*(i//2)
    box(s, cx, cy, cw, ch, fill=LIGHT)
    box(s, cx, cy, Inches(0.14), ch, fill=ORANGE)
    text(s, cx+Inches(0.35), cy+Inches(0.25), cw-Inches(0.6), Inches(0.6), t_, size=19, color=NAVY, bold=True)
    text(s, cx+Inches(0.35), cy+Inches(1.0), cw-Inches(0.6), Inches(1.2), b_, size=15, color=GRAY)
footer(s)

# ================= Slide 3 — UTTEC 소개 =================
s = add_slide(); header(s, 3, "UTTEC — LoRa 무선 제어 8년 실적", "국내 최초 LoRa 조명 콘트롤러 상용화 회사")
rows = [
    ("2018", "국내 최초 LoRa 조명 콘트롤러 상용화 + 클라우드"),
    ("2019", "하나금융 글로벌캠퍼스 — 무선 조명 (사옥 전체)"),
    ("2020", "필로스CC · 광릉CC (한림그룹) — 골프장 LoRa 무선 조명 (코스 전체)"),
    ("2026", "한림용인CC (한림그룹) — 골프장 LoRa 무선 수조 제어 (14노드, 설치 완료)"),
]
y = Inches(1.7)
for yr, desc in rows:
    box(s, Inches(0.6), y, Inches(1.5), Inches(0.85), fill=NAVY)
    text(s, Inches(0.6), y, Inches(1.5), Inches(0.85), yr, size=22, color=WHITE, bold=True, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    box(s, Inches(2.25), y, Inches(10.4), Inches(0.85), fill=LIGHT)
    text(s, Inches(2.5), y, Inches(10.0), Inches(0.85), desc, size=16, color=NAVY, anchor=MSO_ANCHOR.MIDDLE)
    y += Inches(1.05)
text(s, Inches(0.6), Inches(6.5), Inches(12), Inches(0.5),
     "→ 한림그룹 반복 거래(필로스·광릉·용인) + 조명에서 수조로 응용 확장 = 검증된 무선 제어 인프라", size=15, color=ORANGE, bold=True)
footer(s)

# ================= Slide 4 — 2대 시스템 개요 (도식) =================
s = add_slide(); header(s, 4, "하나의 LoRa 무선 인프라, 두 개의 응용", "같은 통신·제어 플랫폼 위에서 수조와 조명을 함께 운용")
# 중앙 허브
hub = box(s, Inches(5.4), Inches(2.9), Inches(2.5), Inches(1.1), fill=NAVY)
text(s, Inches(5.4), Inches(2.9), Inches(2.5), Inches(1.1), "LoRa 무선\n제어 인프라", size=16, color=WHITE, bold=True, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
# 좌 수조
box(s, Inches(0.8), Inches(2.6), Inches(3.6), Inches(1.7), fill=LIGHT)
text(s, Inches(0.8), Inches(2.75), Inches(3.6), Inches(0.5), "시스템 A — 수조 관리", size=17, color=ORANGE, bold=True, align=PP_ALIGN.CENTER)
text(s, Inches(1.0), Inches(3.3), Inches(3.2), Inches(0.9), "자동 수위 판단 → 펌프 제어\n→ 원격 모니터링", size=14, color=GRAY, align=PP_ALIGN.CENTER)
# 우 조명
box(s, Inches(8.9), Inches(2.6), Inches(3.6), Inches(1.7), fill=LIGHT)
text(s, Inches(8.9), Inches(2.75), Inches(3.6), Inches(0.5), "시스템 B — 야간 조명", size=17, color=ORANGE, bold=True, align=PP_ALIGN.CENTER)
text(s, Inches(9.1), Inches(3.3), Inches(3.2), Inches(0.9), "팀 진행 추적 → 미사용 홀\n순차 소등", size=14, color=GRAY, align=PP_ALIGN.CENTER)
# 연결선
for (x1) in [Inches(4.4), Inches(7.9)]:
    cn = box(s, x1, Inches(3.42), Inches(1.0), Pt(3), fill=ORANGE)
text(s, Inches(0.6), Inches(5.0), Inches(12), Inches(1.2),
     "공유 플랫폼: 단일 통합 펌웨어 · 스마트폰 현장 설정(프로비저닝) · 지형 우회 다중 중계 · 원격 대시보드\n→ 한 번 검증한 무선 스택을 두 응용에 재사용 → 개발·유지보수 비용 최소화",
     size=15, color=NAVY)
footer(s)

# ================= Slide 5 — 시스템 A 수조 =================
s = add_slide(); header(s, 5, "시스템 A — 무선 수조 자동 관리", "한림용인CC 14노드 현장 실증 완료")
bullets(s, Inches(0.6), Inches(1.55), Inches(7.3), Inches(5),[
    ("각 센서가 목표 수위를 스스로 판단 → 펌프 자동 ON/OFF", 0, NAVY, True),
    ("본부 서버 없이 현장 노드끼리 자율 운전(폐루프)", 1, GRAY, False),
    ("통신 두절에도 안전(펌프 OFF fail-safe)", 1, GRAY, False),
    ("원거리 수원·펌프·고지대 탱크를 무선으로 연결", 0, NAVY, True),
    ("수위·배터리 실시간 원격 모니터링(대시보드)", 1, GRAY, False),
    ("저전력·솔라 전원 → 전기 없는 외곽 pond까지 대응", 0, NAVY, True),
    ("실증: 한림용인CC 14노드 2단 cascade(연못→저수조→고수조→관수)", 0, GREEN, True),
])
box(s, Inches(8.2), Inches(1.6), Inches(4.5), Inches(4.7), fill=LIGHT)
text(s, Inches(8.2), Inches(1.75), Inches(4.5), Inches(0.5), "자동 급수 흐름", size=16, color=ORANGE, bold=True, align=PP_ALIGN.CENTER)
flow = ["연못(수원)", "펌프 (원거리)", "중계기", "고가 수조", "관수 / 스프링클러"]
fy = Inches(2.4)
for i, f in enumerate(flow):
    box(s, Inches(9.0), fy, Inches(2.9), Inches(0.55), fill=WHITE, line=NAVY)
    text(s, Inches(9.0), fy, Inches(2.9), Inches(0.55), f, size=13, color=NAVY, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    fy += Inches(0.78)
footer(s)

# ================= Slide 6 — 수위 변화 추적 · 누수 조기 진단 (v2, 2026-07-11 추가) =================
s = add_slide(); header(s, 6, "수위 변화 추적 — 배관 누수 조기 진단",
                        "'채우고 비우기'를 넘어, 이상을 진단하는 예방정비 (2026-07-11 한림용인CC 현장 발견)")
# 좌: 기존 부표(오뚜기) 방식
box(s, Inches(0.6), Inches(1.55), Inches(5.9), Inches(3.4), fill=LIGHT)
box(s, Inches(0.6), Inches(1.55), Inches(5.9), Inches(0.62), fill=GRAY)
text(s, Inches(0.6), Inches(1.55), Inches(5.9), Inches(0.62), "기존 부표(오뚜기) 방식",
     size=17, color=WHITE, bold=True, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
bullets(s, Inches(0.95), Inches(2.5), Inches(5.3), Inches(2.3), [
    ("수조의 '가득 / 바닥' 두 지점만 ON·OFF 감지", 0, NAVY, True),
    ("수위가 왜·얼마나 줄어드는지 변화 과정은 못 봄", 1, GRAY, False),
    ("누수·이상은 바닥에 닿을 때까지 발견 불가", 0, GRAY, True),
], size=16, gap=12)
# 우: UTTEC 수위 변화 추적
box(s, Inches(6.85), Inches(1.55), Inches(5.85), Inches(3.4), fill=LIGHT)
box(s, Inches(6.85), Inches(1.55), Inches(5.85), Inches(0.62), fill=GREEN)
text(s, Inches(6.85), Inches(1.55), Inches(5.85), Inches(0.62), "UTTEC 수위 변화 추적",
     size=17, color=WHITE, bold=True, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
bullets(s, Inches(7.2), Inches(2.5), Inches(5.2), Inches(1.4), [
    ("수위를 연속 측정 → 시간에 따른 변화(추세) 관측", 0, NAVY, True),
    ("관수(스프링클러) 미작동 중 수위 하강 = 배관 누수 신호", 0, ORANGE, True),
], size=16, gap=10)
# 우측 하단 미니 추이 그래프 (수위 하강 = 누수)
box(s, Inches(7.15), Inches(3.9), Inches(5.25), Inches(0.9), fill=WHITE, line=GREEN)
text(s, Inches(7.3), Inches(3.94), Inches(2.2), Inches(0.32), "수위 추이", size=10, color=GRAY)
_pts = [(7.55,4.5),(8.45,4.55),(9.35,4.63),(10.25,4.7),(11.0,4.76)]
for i in range(len(_pts)-1):
    c = s.shapes.add_connector(MSO_CONNECTOR.STRAIGHT,
        Inches(_pts[i][0]), Inches(_pts[i][1]), Inches(_pts[i+1][0]), Inches(_pts[i+1][1]))
    c.line.color.rgb = GREEN; c.line.width = Pt(2.5)
text(s, Inches(11.05), Inches(4.4), Inches(1.3), Inches(0.4), "누수 ↓", size=13, color=RED, bold=True, anchor=MSO_ANCHOR.MIDDLE)
# 하단 배너 — 현장 실증 + 원인 진단 프레임
box(s, Inches(0.6), Inches(5.2), Inches(12.1), Inches(1.5), fill=NAVY)
text(s, Inches(0.9), Inches(5.32), Inches(11.5), Inches(0.7),
     "현장 실증: 한림용인CC에서 스프링클러 미작동 중 수위 하강을 포착 → 배관 누수 확인.",
     size=16, color=WHITE, bold=True)
text(s, Inches(0.9), Inches(5.95), Inches(11.5), Inches(0.7),
     "경쟁의 '수위 유지 제어 = 증상 대응' vs UTTEC '수위 변화 추적 = 원인 진단(예방정비)'"
     "   ※ 누수 자동 알림은 SW 로드맵",
     size=14, color=ORANGE, bold=True)
footer(s)

# ================= Slide 7 — 시스템 B 조명 =================
s = add_slide(); header(s, 7, "시스템 B — 야간 조명 순차 제어", "무선 조명 8년 상용화 실적 + 순차 소등 제어 설계 적용")
bullets(s, Inches(0.6), Inches(1.55), Inches(7.3), Inches(5),[
    ("야간 개장 시 마지막 팀 진행에 따라 미사용 홀 조명을 순차 소등", 0, NAVY, True),
    ("사용 중인 홀만 켜두고 지나간 홀은 자동 OFF", 1, GRAY, False),
    ("효과 ① 야간 조명 전력비 절감", 0, ORANGE, True),
    ("효과 ② 담당자가 직접 스위치 끄러 가는 불편 제거", 0, ORANGE, True),
    ("효과 ③ 야간 이동 중 사고 최소화", 0, ORANGE, True),
    ("UTTEC 무선 조명 제어는 2018년 국내 최초 상용화 후 다현장 운영", 0, GREEN, True),
    ("순차 소등(진행 추적) 로직은 본 제안의 신규 설계 적용 항목", 1, GRAY, False),
])
box(s, Inches(8.2), Inches(1.6), Inches(4.5), Inches(4.7), fill=LIGHT)
text(s, Inches(8.2), Inches(1.75), Inches(4.5), Inches(0.5), "순차 소등 개념", size=16, color=ORANGE, bold=True, align=PP_ALIGN.CENTER)
holes = [("1~9홀","지나감 · OFF", GRAY), ("10~12홀","진행 중 · ON", GREEN), ("13~18홀","대기 · ON", NAVY)]
hy = Inches(2.5)
for name, st, col in holes:
    box(s, Inches(9.0), hy, Inches(2.9), Inches(0.9), fill=WHITE, line=col)
    text(s, Inches(9.15), hy+Inches(0.1), Inches(2.6), Inches(0.4), name, size=14, color=NAVY, bold=True)
    text(s, Inches(9.15), hy+Inches(0.45), Inches(2.6), Inches(0.4), st, size=12, color=col, bold=True)
    hy += Inches(1.1)
footer(s)

# ================= Slide 7 — 안전 가치 =================
s = add_slide(); header(s, 8, "★ 안전 — 인적 사고 예방", "전력비·편의를 넘어서는 핵심 가치")
box(s, Inches(0.6), Inches(1.6), Inches(12.1), Inches(1.5), fill=LIGHT)
text(s, Inches(0.9), Inches(1.75), Inches(11.5), Inches(1.2),
     "골프장 지형상 사고가 나기 쉬운 위치(경사지·외곽·수변·펌프실)의 차단기·스위치를\n야간에 직접 조작하지 않고 원격으로 제어 → 이동·조작 중 인적 사고를 근본 예방",
     size=18, color=NAVY, bold=True, anchor=MSO_ANCHOR.MIDDLE)
pairs = [
    ("야간 조명 소등 이동", "→ 원격 순차 소등, 야간 이동 불요"),
    ("경사지 차단기 수동 조작", "→ 원격 제어, 위험 위치 접근 최소화"),
    ("우천·야간 펌프 조작", "→ 자동 제어, 감전·낙상 위험 회피"),
]
py = Inches(3.5)
for a,b in pairs:
    box(s, Inches(0.9), py, Inches(5.3), Inches(0.85), fill=WHITE, line=GRAY)
    text(s, Inches(1.1), py, Inches(5.0), Inches(0.85), a, size=15, color=GRAY, anchor=MSO_ANCHOR.MIDDLE)
    box(s, Inches(6.5), py, Inches(6.2), Inches(0.85), fill=WHITE, line=GREEN)
    text(s, Inches(6.7), py, Inches(5.9), Inches(0.85), b, size=15, color=GREEN, bold=True, anchor=MSO_ANCHOR.MIDDLE)
    py += Inches(1.0)
footer(s)

# ================= Slide 8 — 재고 AS킷 =================
s = add_slide(); header(s, 9, "재고 AS킷 — 긴급 고장 현장 1탭 복구", "1홀 = 1 표준 케이스 + 공용 예비품")
bullets(s, Inches(0.6), Inches(1.6), Inches(6.1), Inches(5),[
    ("모든 제어점을 동일한 표준 케이스로 구성", 0, NAVY, True),
    ("전 케이스가 완전히 같은 펌웨어 탑재", 1, GRAY, False),
    ("설치·교체 시 스마트폰 앱으로 담당 위치만 설정", 1, GRAY, False),
    ("예비 케이스 = 어느 위치든 될 수 있는 공용 재고", 0, NAVY, True),
    ("고장 시 재고 케이스로 현장 즉시 교체(출장 1회·수십 분)", 0, ORANGE, True),
    ("단일 고장이 여러 홀로 번지지 않음(1홀만 영향)", 1, GRAY, False),
])
box(s, Inches(7.0), Inches(1.6), Inches(5.7), Inches(4.7), fill=LIGHT)
text(s, Inches(7.0), Inches(1.75), Inches(5.7), Inches(0.5), "긴급 AS 절차", size=16, color=ORANGE, bold=True, align=PP_ALIGN.CENTER)
steps = ["① 재고함에서 예비 케이스","② 앱으로 해당 위치 선택·설정","③ 고장 케이스 교체·결선","④ 동작 확인 → 완료"]
sy = Inches(2.45)
for st in steps:
    box(s, Inches(7.4), sy, Inches(4.9), Inches(0.7), fill=WHITE, line=NAVY)
    text(s, Inches(7.6), sy, Inches(4.6), Inches(0.7), st, size=14, color=NAVY, anchor=MSO_ANCHOR.MIDDLE)
    sy += Inches(0.92)
footer(s)

# ================= Slide 9 — 레퍼런스 =================
s = add_slide(); header(s, 10, "레퍼런스 — 운영 중인 현장", "조명 · 수조, 한림그룹 다현장")
refs = [
    ("필로스CC", "골프장 무선 조명", "2020 · 코스 전체"),
    ("광릉CC", "골프장 무선 조명", "2020 · 코스 전체"),
    ("하나금융 글로벌캠퍼스", "사옥 무선 조명", "2019"),
    ("한림용인CC", "골프장 무선 수조 제어", "2026 · 14노드 설치 완료"),
]
x0=Inches(0.6); y0=Inches(1.7); cw=Inches(6.0); ch=Inches(2.2); gx=Inches(0.3); gy=Inches(0.35)
for i,(n,sy_,d) in enumerate(refs):
    cx=x0+(cw+gx)*(i%2); cy=y0+(ch+gy)*(i//2)
    box(s, cx, cy, cw, ch, fill=LIGHT)
    box(s, cx, cy, cw, Inches(0.6), fill=NAVY)
    text(s, cx+Inches(0.25), cy, cw-Inches(0.5), Inches(0.6), n, size=18, color=WHITE, bold=True, anchor=MSO_ANCHOR.MIDDLE)
    text(s, cx+Inches(0.25), cy+Inches(0.8), cw-Inches(0.5), Inches(0.5), sy_, size=16, color=ORANGE, bold=True)
    text(s, cx+Inches(0.25), cy+Inches(1.4), cw-Inches(0.5), Inches(0.5), d, size=14, color=GRAY)
footer(s)

# ================= Slide 10 — 도입 효과 =================
s = add_slide(); header(s, 11, "도입 효과", "현행 대비 개선")
data = [
    ("항목", "현행", "LoRa 통합 제어 도입"),
    ("야간 조명 전력", "미사용 홀도 점등 유지", "진행 따라 순차 소등 → 절감"),
    ("인력·이동", "야간 수동 소등·수위 점검 이동", "무인 자동 운전"),
    ("긴급 AS", "노드별 펌웨어·재방문", "재고 케이스 현장 1탭 복구"),
    ("안전", "위험 위치 야간 수동 조작", "원격 제어 → 사고 예방"),
]
tblw = Inches(12.1); x=Inches(0.6); y=Inches(1.7)
colw=[Inches(3.1), Inches(4.5), Inches(4.5)]; rh=Inches(0.9)
for ri,row in enumerate(data):
    cx=x
    for ci,cell in enumerate(row):
        fill = NAVY if ri==0 else (LIGHT if ri%2 else WHITE)
        col  = WHITE if ri==0 else (NAVY if ci==0 else GRAY)
        b = box(s, cx, y, colw[ci], rh, fill=fill, line=RGBColor(0xD0,0xD7,0xDE))
        bd = (ri==0) or (ci==0) or (ci==2)
        c2 = GREEN if (ri>0 and ci==2) else col
        text(s, cx+Inches(0.15), y, colw[ci]-Inches(0.3), rh, cell, size=14, color=c2, bold=bd, anchor=MSO_ANCHOR.MIDDLE)
        cx += colw[ci]
    y += rh
footer(s)

# ================= Slide 11 — 현장 적용 프로세스 =================
s = add_slide(); header(s, 12, "현장 적용 프로세스", "검증된 4단계 도입 절차")
steps = [
    ("1. 현장 RF 서베이", "지형·차폐 실측\n중계 배치·통신 설계"),
    ("2. 시스템 설계", "제어점·노드·중계 수량\n규모별 견적 확정"),
    ("3. 시공·시운전", "표준 케이스 설치\n앱 프로비저닝·통합 시험"),
    ("4. 교육·SOP", "재고 설정·교체 절차\n운영자 인수인계"),
]
x=Inches(0.6); w=Inches(2.95); gap=Inches(0.15); y=Inches(2.4); h=Inches(2.2)
for i,(t_,b_) in enumerate(steps):
    cx = x + (w+gap)*i
    box(s, cx, y, w, h, fill=LIGHT)
    box(s, cx, y, w, Inches(0.65), fill=ORANGE)
    text(s, cx, y, w, Inches(0.65), t_, size=15, color=WHITE, bold=True, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    text(s, cx+Inches(0.2), y+Inches(0.8), w-Inches(0.4), Inches(1.3), b_, size=13, color=GRAY, align=PP_ALIGN.CENTER)
    if i < 3:
        text(s, cx+w-Inches(0.02), y+Inches(0.7), Inches(0.2), Inches(0.6), "▶", size=16, color=ORANGE, align=PP_ALIGN.CENTER)
text(s, Inches(0.6), Inches(5.1), Inches(12), Inches(0.6),
     "무료 현장 진단 / RF 서베이부터 시작합니다 — 규모·환경을 확인한 뒤 정확한 견적을 드립니다.", size=15, color=NAVY, bold=True)
footer(s)

# ================= Slide 12 — 견적 프레임 =================
s = add_slide(); header(s, 13, "견적 프레임 — 규모별 3단계", "구체 단가는 현장 진단 후 제시")
data = [
    ("등급", "구성", "예상 수명 / 하자보증", "적용"),
    ("초경제형", "경제형 컨트롤러·모듈", "5~7년 / 6개월", "소규모·예산 제약"),
    ("표준", "산업용 LoRa·솔라 중계", "8~10년 / 1년", "중규모·NLOS 다수"),
    ("산업급", "임베디드·풀 제어반·HMI", "10~15년 / 1~2년", "대형·무인 안정성"),
]
x=Inches(0.6); y=Inches(1.7); colw=[Inches(2.2),Inches(4.2),Inches(3.4),Inches(2.3)]; rh=Inches(0.95)
for ri,row in enumerate(data):
    cx=x
    for ci,cell in enumerate(row):
        fill = NAVY if ri==0 else (LIGHT if ri%2 else WHITE)
        col  = WHITE if ri==0 else (ORANGE if ci==0 else NAVY if ci==0 else GRAY)
        box(s, cx, y, colw[ci], rh, fill=fill, line=RGBColor(0xD0,0xD7,0xDE))
        c = WHITE if ri==0 else (ORANGE if ci==0 else GRAY)
        text(s, cx+Inches(0.15), y, colw[ci]-Inches(0.3), rh, cell, size=14, color=c, bold=(ri==0 or ci==0), anchor=MSO_ANCHOR.MIDDLE)
        cx += colw[ci]
    y += rh
box(s, Inches(0.6), Inches(5.8), Inches(12.1), Inches(1.0), fill=LIGHT)
text(s, Inches(0.85), Inches(5.9), Inches(11.6), Inches(0.85),
     "총액은 제어점(홀/노드) 수에 비례하며, 게이트웨이·관제·중계·용역은 공통 고정비입니다.\n2회차부터 통합 펌웨어 재사용으로 개발비가 절감됩니다. — 상세 견적은 현장 진단 후 제시.",
     size=13, color=NAVY, anchor=MSO_ANCHOR.MIDDLE)
footer(s)

# ================= Slide 13 — 경제성 개요 =================
s = add_slide(); header(s, 14, "경제성 — 투자 대비 절감", "도입 기준 비용 + 공개 요금(2025)으로 회수 계산")
box(s, Inches(0.6), Inches(1.6), Inches(6.0), Inches(4.7), fill=LIGHT)
text(s, Inches(0.6), Inches(1.75), Inches(6.0), Inches(0.5), "도입 기준 비용", size=18, color=ORANGE, bold=True, align=PP_ALIGN.CENTER)
bullets(s, Inches(0.95), Inches(2.4), Inches(5.4), Inches(3.6), [
    ("수조 1 set (펌프 제어 + 센서) — 150만원", 0, NAVY, True),
    ("조명 1 홀 제어 — 50만원", 0, NAVY, True),
    ("현장 검토비 — 50만원", 0, GRAY, False),
    ("현장 설치비 — 350만원", 0, GRAY, False),
], size=17, gap=14)
box(s, Inches(6.9), Inches(1.6), Inches(5.8), Inches(4.7), fill=LIGHT)
text(s, Inches(6.9), Inches(1.75), Inches(5.8), Inches(0.5), "절감 계산 요금 기준", size=18, color=ORANGE, bold=True, align=PP_ALIGN.CENTER)
bullets(s, Inches(7.25), Inches(2.4), Inches(5.2), Inches(3.6), [
    ("전기 160 원/kWh (한전 일반용 전력)", 0, NAVY, True),
    ("상수도 1,300 원/㎥ (영업용·군지역 골프장)", 0, NAVY, True),
    ("전국 수돗물 평균 829원/㎥, 군지역 1,017원", 1, GRAY, False),
    ("운영량(부하·시간·인건비)은 현장별 조정", 1, GRAY, False),
], size=16, gap=12)
footer(s)

# ================= Slide 14 — 조명 경제성 =================
s = add_slide(); header(s, 15, "경제성 ① 야간 조명 (18홀 예시)", "순차 소등 = 평균 점등시간 절반 + 수동 소등 인건비 제거")
stat(s, Inches(0.6),  Inches(1.6), Inches(3.9), "투자", "1,300만원", "조명18홀+검토+설치")
stat(s, Inches(4.7),  Inches(1.6), Inches(3.9), "연간 절감", "666만원", "전력 216 + 인건비 450", vcolor=GREEN)
stat(s, Inches(8.8),  Inches(1.6), Inches(3.9), "회수기간", "약 2년", "5년 순이익 2,030만원", vcolor=NAVY)
bullets(s, Inches(0.7), Inches(3.7), Inches(12), Inches(2.6), [
    ("전력 절감 216만원/년 = 36kW × 2.5h × 150일 × 160원", 0, NAVY, False),
    ("마지막 팀 5시간 소요 → 순차 소등 시 평균 2.5시간만 점등", 1, GRAY, False),
    ("인건비 절감 450만원/년 = 야간 수동 소등 1.5h × 150일 × 20,000원", 0, NAVY, False),
], size=16, gap=12)
text(s, Inches(0.7), Inches(6.4), Inches(12), Inches(0.5),
     "※ 야간 개장일수·조명 부하·인건비는 현장별 조정 (위는 보수적 예시)", size=12, color=GRAY)
footer(s)

# ================= Slide 15 — 수조 경제성 =================
s = add_slide(); header(s, 16, "경제성 ② 수조 관리 (2 set 예시)", "수위 순회 인건비 제거 + 비싼 긴급 상수도 회피")
stat(s, Inches(0.6),  Inches(1.6), Inches(3.9), "투자", "700만원", "수조2set+검토+설치")
stat(s, Inches(4.7),  Inches(1.6), Inches(3.9), "연간 절감", "795만원", "인건비 600 + 상수도 195", vcolor=GREEN)
stat(s, Inches(8.8),  Inches(1.6), Inches(3.9), "회수기간", "약 11개월", "5년 순이익 3,275만원", vcolor=NAVY)
bullets(s, Inches(0.7), Inches(3.7), Inches(12), Inches(2.6), [
    ("인건비 절감 600만원/년 = 수위 순회 점검 1.0h × 300일 × 20,000원", 0, NAVY, False),
    ("긴급 상수도 회피 195만원/년 = 1,500㎥ × 1,300원", 0, NAVY, False),
    ("물 부족 시 비싼 상수도 급수 대신, 자동 관리로 미리 확보", 1, GRAY, False),
], size=16, gap=12)
text(s, Inches(0.7), Inches(6.4), Inches(12), Inches(0.5),
     "※ 통합 도입(조명+수조) 시 검토·설치 1회 부담 → 회수 약 1.1년 · 5년 순이익 5,705만원", size=12, color=ORANGE, bold=True)
footer(s)

# ================= Slide 16 — 정직·신뢰 =================
s = add_slide(); header(s, 17, "정직하게 관리하는 기술 변수", "숨기지 않고, 설계·시공으로 해결합니다")
pairs = [
    ("무선 채널 용량", "제어점이 많은 대형 현장은 규모 상한·중계·시분할을 설계에 반영"),
    ("통신 커버리지", "지형 차폐 구간은 현장 RF 서베이·안테나 배치로 사전 확보"),
    ("현장 결선·개체 편차", "표준 케이스 + 설치 SOP + 점검 도구로 시공 품질 보증"),
]
y=Inches(1.8)
for a,b in pairs:
    box(s, Inches(0.6), y, Inches(3.6), Inches(1.2), fill=NAVY)
    text(s, Inches(0.8), y, Inches(3.2), Inches(1.2), a, size=16, color=WHITE, bold=True, anchor=MSO_ANCHOR.MIDDLE)
    box(s, Inches(4.4), y, Inches(8.3), Inches(1.2), fill=LIGHT)
    text(s, Inches(4.65), y, Inches(7.9), Inches(1.2), b, size=15, color=NAVY, anchor=MSO_ANCHOR.MIDDLE)
    y += Inches(1.45)
text(s, Inches(0.6), Inches(6.5), Inches(12), Inches(0.5),
     "→ 변수를 아는 것과 도구·SOP로 관리하는 것이 UTTEC의 신뢰입니다.", size=15, color=ORANGE, bold=True)
footer(s)

# ================= Slide 14 — CTA =================
s = add_slide()
box(s, 0, 0, SW, SH, fill=NAVY)
box(s, 0, Inches(4.4), SW, Pt(5), fill=ORANGE)
text(s, Inches(0.9), Inches(1.6), Inches(11.5), Inches(1.2),
     "무료 현장 진단으로 시작하세요", size=38, color=WHITE, bold=True)
text(s, Inches(0.9), Inches(3.0), Inches(11.5), Inches(1.2),
     "귀 골프장의 규모·지형을 확인한 뒤,\n수조·조명 통합 제어 도입 계획과 견적을 제시해 드립니다.",
     size=20, color=RGBColor(0xE2,0xE8,0xF0))
text(s, Inches(0.9), Inches(5.1), Inches(11.5), Inches(0.6),
     "UTTEC  ·  국내 최초 LoRa 무선 조명 상용화(2018)", size=16, color=ORANGE, bold=True)
text(s, Inches(0.9), Inches(5.9), Inches(11.5), Inches(0.6),
     "문의  ·  ihong9059@gmail.com  ·  010-2401-9059", size=16, color=RGBColor(0xCB,0xD5,0xE1))

out = r"C:\todo\today\영업\LoRa원거리제어\제안서\UTTEC_LoRa무선제어_제안서_범용_v2_20260711.pptx"
prs.save(out)
print("saved:", out, "slides:", len(prs.slides._sldIdLst))

# -*- coding: utf-8 -*-
"""UTTEC 미래창의 아카데미 3차 — 수업용 강의 슬라이드 공통 디자인 시스템.
python-pptx 기반. 회차별 build_0N.py에서 import 하여 재사용한다.
대상: 비전공자 초급(SME 임직원·예비창업자) → 큰 글씨·시각 중심·한 슬라이드 한 메시지.
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn
import copy

# ---- 디자인 토큰 ----------------------------------------------------------
FONT = "맑은 고딕"

NAVY   = RGBColor(0x1B, 0x3A, 0x5B)   # 주색 — 제조·신뢰
TEAL   = RGBColor(0x2A, 0x9D, 0x8F)   # 보조 — 이론
ORANGE = RGBColor(0xE7, 0x6F, 0x51)   # 강조 — 실습/체험
YELLOW = RGBColor(0xF4, 0xC4, 0x30)   # 현업노트 제안 배지
LIGHT  = RGBColor(0xF4, 0xF6, 0xF8)   # 밝은 배경
INK    = RGBColor(0x22, 0x2B, 0x35)   # 본문 텍스트
GRAY   = RGBColor(0x8A, 0x97, 0xA3)   # 보조 텍스트
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
CARD   = RGBColor(0xFF, 0xFF, 0xFF)
SOFT   = RGBColor(0xE7, 0xEC, 0xF1)   # 카드 경계/연한 박스
PURPLE = RGBColor(0x6A, 0x4C, 0x93)   # 위키/AI 트랙
PLUM   = RGBColor(0x8E, 0x6F, 0xB8)   # 위키 보조

# 카테고리 칩 색상
CAT = {
    "OT":   NAVY,
    "이론": TEAL,
    "실습": ORANGE,
    "체험": ORANGE,
    "현업노트": YELLOW,
    "위키":  PURPLE,
    "AI활용": PURPLE,
    "미션":  ORANGE,
    "정리": NAVY,
}

EMU_W = Inches(13.333)
EMU_H = Inches(7.5)

PROGRAM = "AI·IoT 스마트팩토리 실습 과정 · 광주"


def new_deck():
    prs = Presentation()
    prs.slide_width = EMU_W
    prs.slide_height = EMU_H
    return prs


def _blank(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])


def blank(prs):
    return _blank(prs)


def _set_fill(shape, color):
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()


def _no_autosize(tf):
    tf.word_wrap = True


def _ea_font(run, name=FONT):
    """한글(동아시아) 폰트까지 강제 지정."""
    run.font.name = name
    rPr = run._r.get_or_add_rPr()
    for tag in ("a:latin", "a:ea", "a:cs"):
        el = rPr.find(qn(tag))
        if el is None:
            el = rPr.makeelement(qn(tag), {})
            rPr.append(el)
        el.set("typeface", name)


def rect(slide, x, y, w, h, color, shape=MSO_SHAPE.RECTANGLE, line=None, line_w=None):
    sp = slide.shapes.add_shape(shape, x, y, w, h)
    sp.shadow.inherit = False
    if color is None:
        sp.fill.background()
    else:
        _set_fill(sp, color)
    if line is not None:
        sp.line.color.rgb = line
        sp.line.width = line_w or Pt(1)
    else:
        sp.line.fill.background()
    return sp


def text(slide, x, y, w, h, runs, size=18, color=INK, bold=False, align=PP_ALIGN.LEFT,
         anchor=MSO_ANCHOR.TOP, line_spacing=1.08, space_after=4, font=FONT):
    """runs: str 또는 [(텍스트, {옵션}) ...] 또는 [[run,run],[run,...]] (문단 리스트)."""
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    _no_autosize(tf)
    tf.vertical_anchor = anchor
    # 문단 정규화
    if isinstance(runs, str):
        paras = [[(runs, {})]]
    elif runs and isinstance(runs[0], tuple):
        paras = [runs]              # 한 문단, 여러 run
    else:
        paras = runs                # 여러 문단
    first = True
    for para in paras:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.alignment = align
        if line_spacing:
            p.line_spacing = line_spacing
        if space_after is not None:
            p.space_after = Pt(space_after)
        if isinstance(para, tuple):
            para = [para]
        for seg, opt in para:
            r = p.add_run()
            r.text = seg
            r.font.size = Pt(opt.get("size", size))
            r.font.bold = opt.get("bold", bold)
            r.font.color.rgb = opt.get("color", color)
            _ea_font(r, opt.get("font", font))
    return tb


def chip(slide, x, y, label, color=None, text_color=WHITE, w=None):
    color = color or CAT.get(label, NAVY)
    w = w or Inches(0.16 * len(label) + 0.5)
    sp = rect(slide, x, y, w, Inches(0.42), color, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    tf = sp.text_frame
    _no_autosize(tf)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = label
    r.font.size = Pt(14); r.font.bold = True; r.font.color.rgb = text_color
    _ea_font(r)
    return sp


def footer(slide, prs, page, total):
    # 하단 얇은 라인 + 과정명 + 페이지
    rect(slide, Inches(0.6), Inches(7.04), Inches(12.13), Pt(1.2), SOFT)
    text(slide, Inches(0.6), Inches(7.08), Inches(9), Inches(0.3),
         [("UTTEC  ·  ", {"bold": True, "color": NAVY, "size": 10}),
          (PROGRAM, {"color": GRAY, "size": 10})], size=10)
    text(slide, Inches(11.4), Inches(7.08), Inches(1.33), Inches(0.3),
         f"{page} / {total}", size=10, color=GRAY, align=PP_ALIGN.RIGHT)


# ---- 슬라이드 템플릿 -------------------------------------------------------

def cover(prs, sess_no, title, subtitle, date, instructors, page_total):
    s = _blank(prs)
    rect(s, 0, 0, EMU_W, EMU_H, NAVY)
    # 우하단 장식 블록
    rect(s, Inches(9.6), Inches(4.7), Inches(3.73), Inches(2.8), TEAL)
    rect(s, Inches(9.6), Inches(4.7), Inches(0.18), Inches(2.8), ORANGE)
    # 상단 브랜드
    text(s, Inches(0.85), Inches(0.7), Inches(8), Inches(0.5),
         [("UTTEC", {"bold": True, "color": WHITE, "size": 22}),
          ("   광주 AI·IoT 스마트팩토리 실습", {"color": RGBColor(0xBF,0xD2,0xE0), "size": 18})], size=22)
    text(s, Inches(0.85), Inches(1.25), Inches(11), Inches(0.5),
         PROGRAM, size=15, color=RGBColor(0xBF,0xD2,0xE0))
    # 회차 라벨
    rect(s, Inches(0.85), Inches(2.55), Inches(2.0), Inches(0.8), ORANGE,
         shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    text(s, Inches(0.85), Inches(2.55), Inches(2.0), Inches(0.8),
         f"{sess_no}회차", size=26, color=WHITE, bold=True,
         align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    # 제목
    text(s, Inches(0.85), Inches(3.6), Inches(10.5), Inches(1.6),
         title, size=46, color=WHITE, bold=True, line_spacing=1.05)
    # 부제
    text(s, Inches(0.9), Inches(5.25), Inches(8.3), Inches(0.6),
         subtitle, size=18, color=RGBColor(0xD9,0xE4,0xEC))
    # 슬로건
    text(s, Inches(0.9), Inches(6.02), Inches(9), Inches(0.4),
         "百見不如一習   ·   직접 익히는 교육", size=13, color=ORANGE, bold=True)
    # 일자·강사
    text(s, Inches(0.9), Inches(6.5), Inches(8), Inches(0.5),
         [(f"📅 {date}", {"color": WHITE, "size": 14, "bold": True}),
          (f"     👤 {instructors}", {"color": RGBColor(0xBF,0xD2,0xE0), "size": 14})], size=14)
    return s


def content(prs, page, total, cat, title, subtitle=None):
    """상단 헤더(칩+제목)만 깔린 콘텐츠 슬라이드 골격 반환."""
    s = _blank(prs)
    rect(s, 0, 0, EMU_W, EMU_H, WHITE)
    # 좌측 색 바
    rect(s, 0, 0, Inches(0.22), EMU_H, CAT.get(cat, NAVY))
    # 헤더
    chip(s, Inches(0.6), Inches(0.55), cat)
    text(s, Inches(0.6), Inches(1.05), Inches(12), Inches(0.9),
         title, size=30, color=NAVY, bold=True)
    if subtitle:
        text(s, Inches(0.62), Inches(1.78), Inches(12), Inches(0.5),
             subtitle, size=15, color=GRAY)
    footer(s, prs, page, total)
    return s


def card(slide, x, y, w, h, color=CARD, line=SOFT, accent=None):
    sp = rect(slide, x, y, w, h, color, shape=MSO_SHAPE.ROUNDED_RECTANGLE,
              line=line, line_w=Pt(1))
    if accent:
        rect(slide, x, y, Inches(0.12), h, accent)
    return sp


def step_card(slide, x, y, w, h, num, head, body, accent=ORANGE):
    card(slide, x, y, w, h, accent=accent)
    # 번호 원
    rect(slide, x + Inches(0.22), y + Inches(0.22), Inches(0.62), Inches(0.62),
         accent, shape=MSO_SHAPE.OVAL)
    text(slide, x + Inches(0.22), y + Inches(0.22), Inches(0.62), Inches(0.62),
         str(num), size=22, color=WHITE, bold=True,
         align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    text(slide, x + Inches(1.0), y + Inches(0.2), w - Inches(1.2), Inches(0.5),
         head, size=17, color=NAVY, bold=True)
    text(slide, x + Inches(1.0), y + Inches(0.72), w - Inches(1.2), h - Inches(0.9),
         body, size=13.5, color=INK, line_spacing=1.12)


def code_box(slide, x, y, w, h, lines, title="실습 코드"):
    rect(slide, x, y, w, h, RGBColor(0x1E, 0x29, 0x33), shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    text(slide, x + Inches(0.3), y + Inches(0.18), w - Inches(0.6), Inches(0.4),
         title, size=12, color=RGBColor(0x7F,0x9A,0xAE), bold=True)
    paras = [[(ln, {"font": "Consolas", "color": RGBColor(0xEA,0xF2,0xF7), "size": 15})] for ln in lines]
    text(slide, x + Inches(0.3), y + Inches(0.62), w - Inches(0.6), h - Inches(0.8),
         paras, size=15, line_spacing=1.25, font="Consolas")


def image_box(slide, x, y, w, h, path):
    """실제 이미지 삽입 (박스 종횡비에 맞춰 생성된 이미지 → 왜곡 없음)."""
    return slide.shapes.add_picture(path, x, y, width=w, height=h)


def photo_fit(slide, x, y, w, h, path, bg=LIGHT, border=True):
    """실물 사진을 박스 안에 종횡비 유지(contain)하며 가운데 배치 → 왜곡 없음.
    박스와 사진의 종횡비가 달라도 letterbox 처리하므로 어떤 사진이든 안전.
    bg/border: 여백(letterbox) 영역을 덮는 둥근 프레임 카드."""
    from PIL import Image
    iw, ih = Image.open(path).size
    box_r = w / h
    img_r = iw / ih
    if img_r > box_r:           # 사진이 더 넓다 → 너비 기준 맞춤
        nw = w; nh = int(round(w * ih / iw))
    else:                       # 사진이 더 높다(또는 같다) → 높이 기준 맞춤
        nh = h; nw = int(round(h * iw / ih))
    if bg is not None:
        rect(slide, x, y, w, h, bg, shape=MSO_SHAPE.ROUNDED_RECTANGLE,
             line=(SOFT if border else None), line_w=Pt(1))
    px = x + (w - nw) // 2
    py = y + (h - nh) // 2
    return slide.shapes.add_picture(path, int(px), int(py), width=int(nw), height=int(nh))


def photo_ph(slide, x, y, w, h, caption):
    """실제 사진 자리 placeholder (회색 박스)."""
    rect(slide, x, y, w, h, LIGHT, shape=MSO_SHAPE.ROUNDED_RECTANGLE,
         line=SOFT, line_w=Pt(1))
    text(slide, x, y + h/2 - Inches(0.45), w, Inches(0.4),
         "🖼", size=26, color=GRAY, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    text(slide, x, y + h/2 + Inches(0.02), w, Inches(0.5),
         caption, size=12, color=GRAY, align=PP_ALIGN.CENTER)


def prompt_card(slide, x, y, w, h, ask, accent=PURPLE):
    """Claude에게 시키는 말풍선 카드."""
    card(slide, x, y, w, h, color=LIGHT, accent=accent)
    text(slide, x+Inches(0.28), y+Inches(0.14), w-Inches(0.5), Inches(0.4),
         "🗣  나 → Claude", size=12, color=accent, bold=True)
    text(slide, x+Inches(0.28), y+Inches(0.56), w-Inches(0.5), h-Inches(0.7),
         ask, size=14, color=INK, line_spacing=1.12)


def wiki_segment(prs, page, total, title, subtitle, activity_head, activity_paras,
                 prompts, why_line, badge=True):
    """회차 공통 '오늘의 위키(Claude 연동)' 세그먼트 슬라이드.
    activity_paras: text()용 문단 리스트 / prompts: Claude 프롬프트 문자열 리스트."""
    s = content(prs, page, total, "위키", title, subtitle)
    if badge:
        proposal_badge(s, Inches(9.95), Inches(0.6))
    # 좌측 — 오늘의 활동
    card(s, Inches(0.6), Inches(2.4), Inches(5.85), Inches(3.95), color=LIGHT, accent=PURPLE)
    text(s, Inches(0.92), Inches(2.62), Inches(5.2), Inches(0.5),
         activity_head, size=17, color=PURPLE, bold=True)
    text(s, Inches(0.92), Inches(3.2), Inches(5.25), Inches(3.0),
         activity_paras, size=15, color=INK, line_spacing=1.3)
    # 우측 — Claude 프롬프트 카드 스택
    px = Inches(6.75); pw = Inches(5.95)
    n = len(prompts)
    ph = Inches(1.18) if n >= 3 else Inches(1.5)
    gap = Inches(0.18)
    y = Inches(2.4)
    for q in prompts:
        prompt_card(s, px, y, pw, ph, q)
        y += ph + gap
    # 하단 — 현업 연결 한 줄
    text(s, Inches(0.6), Inches(6.5), Inches(12.1), Inches(0.4),
         why_line, size=14, color=PURPLE, bold=True, align=PP_ALIGN.CENTER)
    return s


def proposal_badge(slide, x, y):
    """현업노트 강화 = '제안(협의 미확정)' 배지."""
    sp = rect(slide, x, y, Inches(2.45), Inches(0.46), YELLOW,
              shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    text(slide, x, y, Inches(2.45), Inches(0.46),
         "★ 강화 제안 (협의 진행 중)", size=12.5, color=NAVY, bold=True,
         align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
